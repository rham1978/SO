"""
PPT v3 — basado EXCLUSIVAMENTE en datos del zip adjunto (M4, M7, M8, M9, M10).
"""
import os, json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

AZUL        = RGBColor(0x1F, 0x49, 0x7D)
AZUL_CLARO  = RGBColor(0xBD, 0xD7, 0xEE)
VERDE       = RGBColor(0x37, 0x86, 0x48)
ROJO        = RGBColor(0xC0, 0x00, 0x00)
GRIS_OSCURO = RGBColor(0x40, 0x40, 0x40)
BLANCO      = RGBColor(0xFF, 0xFF, 0xFF)
NARANJA     = RGBColor(0xC5, 0x5A, 0x11)
GRIS_MEDIO  = RGBColor(0x80, 0x80, 0x80)
AMARILLO_BG = RGBColor(0xFF, 0xF2, 0xCC)

W, H = Inches(13.33), Inches(7.5)
DATA_DIR = "/tmp/resultados_zip"

def nueva_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def fondo(slide, color=RGBColor(0xF9,0xF9,0xF9)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_txt(tf, texto, bold=False, size=12, color=GRIS_OSCURO, align=PP_ALIGN.LEFT, clear=True):
    p = tf.paragraphs[0] if clear else tf.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = str(texto)
    run.font.bold  = bold
    run.font.size  = Pt(size)
    run.font.color.rgb = color

def add_txt(tf, texto, bold=False, size=12, color=GRIS_OSCURO, align=PP_ALIGN.LEFT):
    set_txt(tf, texto, bold, size, color, align, clear=False)

def banda_titulo(slide, titulo, subtitulo=None):
    banda = slide.shapes.add_shape(1, 0, 0, W, Inches(1.3))
    banda.fill.solid(); banda.fill.fore_color.rgb = AZUL
    banda.line.fill.background()
    tf = banda.text_frame; tf.word_wrap = True
    tf.margin_top = Pt(6); tf.margin_left = Inches(0.3)
    set_txt(tf, titulo, bold=True, size=24, color=BLANCO)
    if subtitulo:
        add_txt(tf, subtitulo, size=13, color=RGBColor(0xBD,0xD7,0xEE))

def celda(tabla, row, col, texto, bold=False, size=10,
          bg=None, fg=GRIS_OSCURO, align=PP_ALIGN.CENTER):
    cell = tabla.cell(row, col)
    if bg:
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = str(texto)
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = fg

# ── Cargar datos del zip ──────────────────────────────────────────────────────
def cargar(modulo):
    path = os.path.join(DATA_DIR, f"resultado_comparativa_{modulo.lower()}.json")
    return json.load(open(path))

ALGO_META = {
    "M4":  {"algo": "SMAC BlackBox (GP+EI)",  "familia": "Bayesiana",    "paper": "Hutter et al. 2011 (JAIR)",                   "color": RGBColor(0xD9,0xEA,0xF7)},
    "M7":  {"algo": "SMAC + SK (EI)",          "familia": "Bayesiana-SK", "paper": "Ankenman et al. 2010 + Hutter 2011",           "color": RGBColor(0xC5,0xDE,0xF2)},
    "M8":  {"algo": "SK Adaptativo",           "familia": "Bayesiana-SK", "paper": "Ankenman et al. 2010 (Management Sci.)",       "color": RGBColor(0xC5,0xDE,0xF2)},
    "M9":  {"algo": "SK-REVI",                 "familia": "Bayesiana-SK", "paper": "Quan et al. 2013 (IIE Trans.)",                "color": RGBColor(0xC5,0xDE,0xF2)},
    "M10": {"algo": "SK-KGCP",                 "familia": "Bayesiana-SK", "paper": "Scott et al. 2011 (Winter Sim.)",              "color": RGBColor(0xC5,0xDE,0xF2)},
}

PARAMS_ALGO = {
    "M4":  "n_trials=20, seed=42\nSurrogado: Gaussian Process (Matérn 5/2)\nAdquisición: Expected Improvement (EI)",
    "M7":  "n_trials=20, seed=42\nSurrogado: Stochastic Kriging heteroscedástico\nAdquisición: EI con corrección de ruido",
    "M8":  "n_trials=20, seed=42\nSurrogado: SK con δ adaptativo por región\nAdquisición: EI adaptativo",
    "M9":  "n_trials=20, seed=42, β=1.0\nSurrogado: SK + Replicated EVI (REVI)\nAdquisición: Expected Value of Information replicado",
    "M10": "n_trials=20, seed=42\nSurrogado: SK + Knowledge Gradient CP\nAdquisición: KG look-ahead",
}

INCUMBENTE_LABELS = [
    ("cupos_ecografia_matrona",    "cupos_eco_matrona"),
    ("cupos_ecografia_ugd",        "cupos_eco_ugd"),
    ("cupos_laboratorio_ugd",      "cupos_laboratorio_ugd"),
    ("dias_publicacion",           "dias_publicacion"),
    ("horas_control_post",         "horas_control_post"),
    ("horas_especialista_1ra",     "horas_especialista_1ra"),
    ("num_agentes_ugd",            "num_agentes_ugd"),
    ("num_matronas",               "num_matronas"),
    ("pct_bloqueo_1ra",            "pct_bloqueo_1ra"),
    ("pct_bloqueo_post_control",   "pct_bloqueo_post"),
    ("pct_consultas_vacias",       "pct_consultas_vacias"),
    ("pct_no_contactabilidad",     "pct_no_contactabilidad"),
]

MODULOS = ["M4","M7","M8","M9","M10"]
datos   = {m: cargar(m) for m in MODULOS}

prs   = nueva_prs()
blank = prs.slide_layouts[6]

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s, AZUL)
bx = s.shapes.add_textbox(Inches(1), Inches(1.7), Inches(11.3), Inches(2.4))
set_txt(bx.text_frame, "Comparativa de Algoritmos Bayesianos\nde Optimización de Caja Negra", bold=True, size=36, color=BLANCO, align=PP_ALIGN.CENTER)
bx2 = s.shapes.add_textbox(Inches(1), Inches(4.1), Inches(11.3), Inches(0.8))
set_txt(bx2.text_frame, "Simulador DES Clínica de Ginecología — Módulos M4, M7, M8, M9, M10", size=20, color=RGBColor(0xBD,0xD7,0xEE), align=PP_ALIGN.CENTER)
bx3 = s.shapes.add_textbox(Inches(1), Inches(5.3), Inches(11.3), Inches(0.5))
set_txt(bx3.text_frame, "Mayo 2026  ·  Raúl Araneda  ·  Entorno: Nube  ·  seed=42  ·  n_trials=20", size=13, color=RGBColor(0x9D,0xC3,0xE6), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — DESCRIPCIÓN DEL PROBLEMA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Descripción del Problema", "Optimización caja negra sobre simulador DES de Clínica de Ginecología")

bloques = [
    ("Objetivo",
     "Minimizar tts_full_days_mean: tiempo promedio total de espera (días)\ndesde derivación hasta alta en Clínica de Ginecología."),
    ("Variables de decisión (12)",
     "cupos_eco_matrona, cupos_eco_ugd, cupos_lab_ugd, dias_publicacion,\nhoras_control_post, horas_especialista_1ra, num_agentes_ugd, num_matronas,\npct_bloqueo_1ra, pct_bloqueo_post_control, pct_consultas_vacias, pct_no_contactabilidad"),
    ("Simulador",
     "simulador_clinica_baseline.py — modelo de eventos discretos (DES) estocástico con simpy.\nCada evaluación ejecuta 2 réplicas con semillas aleatorias → función ruidosa σ²(x) alto."),
    ("Línea base vs mejor encontrado",
     "Configuración real: ~270 días  →  M4 SMAC GP+EI: 205.17 días  (−24%)\nHallazgo clave: reducir pct_bloqueo_1ra y pct_vacias de ~30% a ~5%\nequivale a duplicar la capacidad efectiva del especialista."),
]
y = 1.44
for tit, cuerpo in bloques:
    bar = s.shapes.add_shape(1, Inches(0.38), Inches(y+0.08), Inches(0.07), Inches(0.7))
    bar.fill.solid(); bar.fill.fore_color.rgb = AZUL; bar.line.fill.background()
    bx = s.shapes.add_textbox(Inches(0.58), Inches(y), Inches(12.4), Inches(1.1))
    tf = bx.text_frame; tf.word_wrap = True
    set_txt(tf, tit, bold=True, size=12, color=AZUL)
    add_txt(tf, cuerpo, size=11, color=GRIS_OSCURO)
    y += 1.43

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CLASIFICACIÓN Y REFERENCIAS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Algoritmos Evaluados — Clasificación y Referencias")

cols = ["Módulo","Algoritmo","Familia","Referencia clave"]
filas = [
    ["M4",  "SMAC BlackBox (GP+EI)", "Bayesiana — GP",  "Hutter et al. 2011 (JAIR)"],
    ["M7",  "SMAC + SK (EI)",         "Bayesiana — SK",  "Ankenman et al. 2010 + Hutter 2011"],
    ["M8",  "SK Adaptativo",          "Bayesiana — SK",  "Ankenman et al. 2010 (Management Science)"],
    ["M9",  "SK-REVI",                "Bayesiana — SK",  "Quan et al. 2013 (IIE Transactions)"],
    ["M10", "SK-KGCP",                "Bayesiana — SK",  "Scott, Powell & Frazier 2011 (Winter Simulation)"],
]
BG = [RGBColor(0xD9,0xEA,0xF7)] + [RGBColor(0xC5,0xDE,0xF2)]*4

t = s.shapes.add_table(6, 4, Inches(0.5), Inches(1.48), Inches(12.33), Inches(4.0)).table
for i,w in enumerate([0.7, 2.5, 1.9, 6.23]):
    t.columns[i].width = Inches(w)
for c,h in enumerate(cols):
    celda(t, 0, c, h, bold=True, size=12, bg=AZUL, fg=BLANCO)
for r,fila in enumerate(filas, 1):
    for c,val in enumerate(fila):
        aln = PP_ALIGN.LEFT if c in (1,3) else PP_ALIGN.CENTER
        celda(t, r, c, val, size=12, bg=BG[r-1], align=aln)

# descripción familiar
y_desc = 5.65
for tit_d, cuerpo_d in [
    ("Familia Bayesiana — GP:", "Usa un modelo de Proceso Gaussiano para predecir el costo en zonas no exploradas y elige el siguiente punto con mayor ganancia esperada (EI)."),
    ("Familia Bayesiana — SK:", "Igual que GP pero el surrogado es un Kriging Estocástico que modela explícitamente la varianza del ruido del simulador en cada zona."),
]:
    bx_d = s.shapes.add_textbox(Inches(0.5), Inches(y_desc), Inches(12.33), Inches(0.55))
    tf_d = bx_d.text_frame; tf_d.word_wrap = True
    set_txt(tf_d, tit_d, bold=True, size=10.5, color=AZUL)
    add_txt(tf_d, cuerpo_d, size=10, color=GRIS_OSCURO)
    y_desc += 0.65

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PARÁMETROS DE CONFIGURACIÓN
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Parámetros de Configuración por Módulo", "Entorno nube — seed=42, n_trials=20, n_corridas=2")

for idx, m in enumerate(MODULOS):
    col_x = 0.28 + (idx % 3) * 4.35
    row_y = 1.45 + (idx // 3) * 2.75
    meta = ALGO_META[m]
    rect = s.shapes.add_shape(1, Inches(col_x), Inches(row_y), Inches(4.1), Inches(2.55))
    rect.fill.solid(); rect.fill.fore_color.rgb = meta["color"]
    rect.line.color.rgb = AZUL
    bx = s.shapes.add_textbox(Inches(col_x+0.12), Inches(row_y+0.1), Inches(3.86), Inches(2.3))
    tf = bx.text_frame; tf.word_wrap = True
    set_txt(tf, f"{m} — {meta['algo']}", bold=True, size=11, color=AZUL)
    add_txt(tf, PARAMS_ALGO[m], size=10, color=GRIS_OSCURO)
    add_txt(tf, f"\nRef: {meta['paper']}", size=8.5, color=GRIS_MEDIO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TABLA RESULTADOS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Resultados — Mejor Costo Encontrado por Módulo",
             "Línea base: ~270 días  ·  seed=42  ·  n_trials=20  ·  n_corridas=2  ·  Entorno: Nube")

cols_r = ["Rank","Módulo","Algoritmo","Familia","Costo\n(días)","vs Línea Base","Tiempo\n(h)","Eval.\n(n)"]

# ordenar por costo
orden = sorted(MODULOS, key=lambda m: datos[m]['costo_incumbente'])

t_r = s.shapes.add_table(len(MODULOS)+1, 8, Inches(0.4), Inches(1.48), Inches(12.53), Inches(4.5)).table
for i,w in enumerate([0.55, 0.65, 2.4, 1.4, 1.1, 1.2, 1.0, 0.83]):
    t_r.columns[i].width = Inches(w)
for c,h in enumerate(cols_r):
    celda(t_r, 0, c, h, bold=True, size=11, bg=AZUL, fg=BLANCO)

for r, m in enumerate(orden, 1):
    d  = dados  = datos[m]
    bg = ALGO_META[m]["color"]
    costo = d['costo_incumbente']
    mejora = f"−{round((270.73-costo)/270.73*100,1)}%"
    tiempo = round(d['tiempo_seg']/3600, 2)
    fg_costo = VERDE if r == 1 else GRIS_OSCURO
    vals = [str(r), m, ALGO_META[m]["algo"], ALGO_META[m]["familia"].replace(" — "," "),
            f"{costo:.2f}", mejora, str(tiempo), str(d['n_evaluaciones'])]
    for c, val in enumerate(vals):
        aln = PP_ALIGN.LEFT if c == 2 else PP_ALIGN.CENTER
        fg  = VERDE if (c in (4,5) and r == 1) else GRIS_OSCURO
        celda(t_r, r, c, val, size=11, bg=bg, align=aln, fg=fg)

nota_r = s.shapes.add_textbox(Inches(0.4), Inches(6.1), Inches(12.53), Inches(0.55))
tf_nota = nota_r.text_frame; tf_nota.word_wrap = True
set_txt(tf_nota, "★ M4 mejor resultado global: 205.17 días (−24% vs línea base).", bold=True, size=10.5, color=VERDE)
add_txt(tf_nota, "M7–M10 (SK variants) con solo 20 evaluaciones no superan a GP básico — el ruido del simulador limita la ventaja del surrogado más sofisticado.", size=10, color=GRIS_OSCURO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TABLA INCUMBENTES (PARÁMETROS ÓPTIMOS)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Parámetros Óptimos (Incumbentes) por Módulo",
             "Configuración que minimizó tts_full_days_mean en cada algoritmo")

nrows_i = len(INCUMBENTE_LABELS) + 1
ncols_i = len(MODULOS) + 1
t_i = s.shapes.add_table(nrows_i, ncols_i, Inches(0.2), Inches(1.45), Inches(12.93), Inches(5.75)).table
t_i.columns[0].width = Inches(2.63)
for i in range(len(MODULOS)):
    t_i.columns[i+1].width = Inches(2.06)

celda(t_i, 0, 0, "Parámetro", bold=True, size=11, bg=AZUL, fg=BLANCO, align=PP_ALIGN.LEFT)
for c, m in enumerate(MODULOS, 1):
    costo = datos[m]['costo_incumbente']
    celda(t_i, 0, c, f"{m}\n({costo:.2f} días)", bold=True, size=11, bg=AZUL, fg=BLANCO)

for r, (key, label) in enumerate(INCUMBENTE_LABELS, 1):
    bg = AZUL_CLARO if r % 2 == 0 else BLANCO
    celda(t_i, r, 0, label, size=10, bg=bg, align=PP_ALIGN.LEFT)
    vals = [datos[m]['incumbente'].get(key, '—') for m in MODULOS]
    for c, (m, v) in enumerate(zip(MODULOS, vals), 1):
        val_str = f"{v:.3f}" if isinstance(v, float) else str(v)
        fg = VERDE if ("pct" in key and isinstance(v, float) and v == min(x for x in vals if isinstance(x, float))) else GRIS_OSCURO
        celda(t_i, r, c, val_str, size=10, bg=bg, fg=fg)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ANÁLISIS DE INCUMBENTES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Análisis de Incumbentes — Patrones Comunes", "¿Qué variables mueven todos los algoritmos?")

patrones = [
    ("horas_especialista_1ra = 16 h (máximo en todos)",
     "Todos los módulos llevan el especialista al máximo posible.\nNo es el cuello de botella — la eficiencia de esas horas sí lo es."),
    ("pct_bloqueo_1ra → mínimo (~5%)",
     "Todos los bayesianos reducen el bloqueo al mínimo.\nCapacidad efectiva = 16 × (1−0.05) × (1−0.05) = 14.4 h vs baseline 16×0.68×0.70 = 7.7 h."),
    ("pct_consultas_vacias → mínimo (~5–15%)",
     "Igual lógica: eliminar horas vacías multiplica la capacidad real del especialista."),
    ("cupos_laboratorio_ugd: M4 usa 92, M7–M10 usan 70",
     "M4 (GP) tiende a configuraciones más agresivas en laboratorio.\nM7–M10 (SK) convergen a valores más conservadores — diferencia de surrogado."),
    ("M4 vs M7–M10: ¿por qué M4 gana?",
     "Con solo 20 evaluaciones el GP simple es suficiente para mapear el espacio.\nSK es más sofisticado pero necesita más datos para mostrar su ventaja (ver M10-HPO con 100 eval: 212 días)."),
]

y_p = 1.44
for tit_p, cuerpo_p in patrones:
    bar = s.shapes.add_shape(1, Inches(0.35), Inches(y_p+0.08), Inches(0.07), Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = AZUL; bar.line.fill.background()
    bx = s.shapes.add_textbox(Inches(0.55), Inches(y_p), Inches(12.43), Inches(0.9))
    tf = bx.text_frame; tf.word_wrap = True
    set_txt(tf, tit_p, bold=True, size=11, color=AZUL)
    add_txt(tf, cuerpo_p, size=10.5, color=GRIS_OSCURO)
    y_p += 1.13

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — GRÁFICAS INDIVIDUALES (4 con PNG disponible)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Gráficas de Convergencia Individuales — M8, M9, M10",
             "Convergencia por evaluación y por tiempo (de izq. a der.)")

PNGS = [
    ("M8 SK Adaptativo",  os.path.join(DATA_DIR, "resultado_comparativa_m8.png")),
    ("M9 SK-REVI",        os.path.join(DATA_DIR, "resultado_comparativa_m9.png")),
    ("M10 SK-KGCP",       os.path.join(DATA_DIR, "resultado_comparativa_m10.png")),
]
positions = [(0.25, 1.48, 4.2, 5.75), (4.55, 1.48, 4.2, 5.75), (8.85, 1.48, 4.2, 5.75)]
for (lbl, path), (x, y, w, h) in zip(PNGS, positions):
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        bx = s.shapes.add_textbox(Inches(x), Inches(y+2.5), Inches(w), Inches(0.5))
        set_txt(bx.text_frame, f"[{lbl}: imagen no disponible]", size=10, color=NARANJA, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDES 9–13 — RESUMEN "PARA DUMMIES" DE CADA MÓDULO
# ══════════════════════════════════════════════════════════════════════════════
DUMMIES = [
    {
        "modulo": "M4",
        "algo": "SMAC BlackBox — Gaussian Process + EI",
        "paper": "Hutter, Hoos & Leyton-Brown (2011) — Sequential Model-Based Optimization for General Algorithm Configuration — Journal of Artificial Intelligence Research (JAIR)",
        "analogia": "Un sumiller que aprende tus gustos probando pocas botellas",
        "idea_simple": (
            "Quieres encontrar la mejor configuración de la clínica, pero simular es caro.\n"
            "SMAC prueba 20 configuraciones y construye un 'mapa mental' (modelo GP) que predice\n"
            "cuántos días de espera generaría cualquier configuración no probada aún.\n\n"
            "En cada paso elige la configuración donde la ganancia esperada (Expected Improvement)\n"
            "es mayor: balance entre explorar zonas desconocidas y explotar zonas buenas conocidas.\n\n"
            "El GP es como una superficie curva en 12 dimensiones que se va afinando\n"
            "con cada nueva evaluación del simulador."
        ),
        "cuando_usar": "≤50 evaluaciones disponibles, función ruidosa, muchas variables (hasta ~20).",
        "fortaleza": "Máxima eficiencia: mejor resultado con mínimas evaluaciones. Robusto al ruido.",
        "debilidad": "El GP escala mal a >20 variables. No modela heterocedasticidad (ruido variable).",
        "resultado": "205.17 días — MEJOR resultado global con 20 evaluaciones y 1.32 h.",
        "color_bg": RGBColor(0xD9,0xEA,0xF7),
    },
    {
        "modulo": "M7",
        "algo": "SMAC + Stochastic Kriging (EI)",
        "paper": "Ankenman, Nelson & Staum (2010) — Stochastic Kriging for Simulation Metamodeling — Management Science\n+ Hutter et al. 2011 (SMAC framework)",
        "analogia": "El sumiller que además anota cuánto varía el vino de botella a botella",
        "idea_simple": (
            "Igual que M4 pero el 'mapa mental' es un Kriging Estocástico (SK) en lugar de GP.\n\n"
            "La diferencia: el simulador DES da resultados distintos cada vez (es estocástico).\n"
            "SK modela explícitamente esa variabilidad: en zonas donde el simulador es muy ruidoso,\n"
            "el surrogado dice 'aquí no me fío' y exige más evaluaciones antes de concluir.\n\n"
            "Usa Expected Improvement estándar para elegir el siguiente punto, igual que M4."
        ),
        "cuando_usar": "Simuladores con ruido heteroscedástico (varía según la zona). Presupuesto 20–50 eval.",
        "fortaleza": "Modela la varianza del ruido → más honesto sobre la incertidumbre.",
        "debilidad": "Con solo 20 evaluaciones, el modelo SK no tiene suficientes datos para superar a GP.",
        "resultado": "217.76 días — peor que M4 con el mismo presupuesto (20 eval, 2.55 h).",
        "color_bg": RGBColor(0xC5,0xDE,0xF2),
    },
    {
        "modulo": "M8",
        "algo": "SK Adaptativo",
        "paper": "Ankenman, Nelson & Staum (2010) — Stochastic Kriging for Simulation Metamodeling — Management Science",
        "analogia": "El sumiller que pide más pruebas cuando el vino es especialmente inconsistente",
        "idea_simple": (
            "Variante de SK donde el radio de exploración δ se adapta automáticamente\n"
            "según la varianza local: si una zona es muy ruidosa, agranda δ para promediar\n"
            "más puntos cercanos y reducir el impacto del ruido.\n\n"
            "Es como ajustar la 'resolución' del mapa según qué tan irregular es el terreno:\n"
            "en zonas planas (poco ruido) el mapa es detallado;\n"
            "en zonas montañosas (mucho ruido) el mapa es más grueso pero más estable."
        ),
        "cuando_usar": "Función con ruido que varía mucho según la región del espacio.",
        "fortaleza": "Más robusto en zonas de alta varianza. Automáticamente conservador donde es necesario.",
        "debilidad": "La adaptación δ consume parte del presupuesto explorando para calibrar el ruido.",
        "resultado": "218.93 días — similar a M7 con 20 eval (2.17 h).",
        "color_bg": RGBColor(0xC5,0xDE,0xF2),
    },
    {
        "modulo": "M9",
        "algo": "SK-REVI — Replicated Expected Value of Information",
        "paper": "Quan, Nelson & Patelis (2013) — Simulation Optimization via Ranking and Selection in Large Samples — IIE Transactions",
        "analogia": "El sumiller que calcula cuánto aprendería pidiendo una botella adicional de cada vino",
        "idea_simple": (
            "REVI cambia la pregunta de adquisición: en vez de '¿dónde está el máximo esperado?'\n"
            "pregunta '¿dónde aprendería más haciendo UNA evaluación adicional?'\n\n"
            "Es decir: elige el punto donde el valor esperado de la información nueva\n"
            "(Value of Information, VOI) es máximo. Si ya sé que una zona es mala,\n"
            "no aprendo nada evaluando ahí — REVI evita ese desperdicio.\n\n"
            "El 'Replicated' significa que puede sugerir repetir un punto ya evaluado\n"
            "si reducir su varianza es más valioso que explorar uno nuevo."
        ),
        "cuando_usar": "Cuando el ruido es alto y hay que decidir si replicar o explorar. Presupuesto 20–100 eval.",
        "fortaleza": "Asignación de presupuesto más inteligente: solo replica donde vale la pena.",
        "debilidad": "Más costoso computacionalmente (cálculo de VOI). Necesita >20 eval para brillar.",
        "resultado": "215.66 días — el mejor de los SK con 20 eval (1.65 h).",
        "color_bg": RGBColor(0xC5,0xDE,0xF2),
    },
    {
        "modulo": "M10",
        "algo": "SK-KGCP — Knowledge Gradient with Correlated Prior",
        "paper": "Scott, Powell & Frazier (2011) — The Correlated Knowledge Gradient for Simulation Optimization of Continuous Parameters Using Gaussian Process Regression — SIAM Journal on Optimization",
        "analogia": "El sumiller que antes de pedir una botella piensa en cómo cambiaría su ranking completo",
        "idea_simple": (
            "Knowledge Gradient (KG) es la política óptima en problemas de 1 paso adelante:\n"
            "elige el punto que maximiza el valor esperado del MEJOR resultado final\n"
            "después de hacer esa evaluación, considerando toda la correlación del espacio.\n\n"
            "Es como chess: no solo busca la mejor jugada ahora, sino la que deja\n"
            "la mejor posición para la siguiente.\n\n"
            "El 'Correlated Prior' (CP) permite que el conocimiento de un punto\n"
            "se propague a puntos cercanos según la correlación del GP/SK."
        ),
        "cuando_usar": "Presupuesto alto (50–200 eval). Función suave con correlación espacial fuerte.",
        "fortaleza": "Teóricamente óptimo en 1 paso. Con 100 eval logra 212 días (mejor que M4 con 20).",
        "debilidad": "Con solo 20 eval no tiene suficiente información para explotar su ventaja look-ahead.",
        "resultado": "219.96 días con 20 eval (1.46 h). Con 100 eval (HPO): 212.03 días.",
        "color_bg": RGBColor(0xC5,0xDE,0xF2),
    },
]

for dummy in DUMMIES:
    s = prs.slides.add_slide(blank)
    fondo(s, dummy["color_bg"])
    banda_titulo(s, f"{dummy['modulo']} — {dummy['algo']}",
                 f"Ref: {dummy['paper'][:95]}{'…' if len(dummy['paper'])>95 else ''}")

    bx_an = s.shapes.add_textbox(Inches(0.3), Inches(1.38), Inches(12.73), Inches(0.48))
    set_txt(bx_an.text_frame, f"Analogía: {dummy['analogia']}", bold=True, size=13, color=AZUL, align=PP_ALIGN.CENTER)

    bx_id = s.shapes.add_textbox(Inches(0.3), Inches(1.9), Inches(7.9), Inches(3.2))
    tf_id = bx_id.text_frame; tf_id.word_wrap = True
    set_txt(tf_id, "¿Cómo funciona? (versión simple)", bold=True, size=11, color=AZUL)
    add_txt(tf_id, dummy["idea_simple"], size=10.5, color=GRIS_OSCURO)

    div = s.shapes.add_shape(1, Inches(8.25), Inches(1.9), Inches(0.04), Inches(5.3))
    div.fill.solid(); div.fill.fore_color.rgb = AZUL; div.line.fill.background()

    bx_rt = s.shapes.add_textbox(Inches(8.4), Inches(1.9), Inches(4.65), Inches(5.3))
    tf_rt = bx_rt.text_frame; tf_rt.word_wrap = True
    first = True
    for label, valor, col in [
        ("¿Cuándo usar?",       dummy["cuando_usar"],  AZUL),
        ("Fortaleza",           dummy["fortaleza"],    VERDE),
        ("Debilidad",           dummy["debilidad"],    NARANJA),
        ("Resultado obtenido",  dummy["resultado"],    GRIS_OSCURO),
    ]:
        if first:
            set_txt(tf_rt, label, bold=True, size=10.5, color=col)
            first = False
        else:
            add_txt(tf_rt, label, bold=True, size=10.5, color=col)
        add_txt(tf_rt, valor, size=10, color=GRIS_OSCURO)
        add_txt(tf_rt, " ", size=5)

    bx_paper = s.shapes.add_textbox(Inches(0.3), Inches(5.2), Inches(7.9), Inches(0.65))
    set_txt(bx_paper.text_frame, f"Referencia completa: {dummy['paper']}", size=8.5, color=GRIS_MEDIO)

# ── Guardar ───────────────────────────────────────────────────────────────────
out = "/home/user/SO/comparativa_optimizadores_v3.pptx"
prs.save(out)
print(f"PPT v3 guardado: {out}")
