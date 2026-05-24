"""
PPT v3 — datos EXCLUSIVAMENTE de Raúl (2664058f-todo.7z / Raul_23_05_2026)
Módulos: M4 (SMAC-GP), M10-HPO (SK-KGCP 100 eval), M13 (SPSA), M14 (ALOE)
"""
import os, json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

AZUL        = RGBColor(0x1F, 0x49, 0x7D)
AZUL_CLARO  = RGBColor(0xBD, 0xD7, 0xEE)
VERDE       = RGBColor(0x37, 0x86, 0x48)
ROJO        = RGBColor(0xC0, 0x00, 0x00)
GRIS_OSCURO = RGBColor(0x40, 0x40, 0x40)
BLANCO      = RGBColor(0xFF, 0xFF, 0xFF)
NARANJA     = RGBColor(0xC5, 0x5A, 0x11)
GRIS_MEDIO  = RGBColor(0x80, 0x80, 0x80)
AMARILLO_BG = RGBColor(0xFF, 0xF2, 0xCC)
ROJO_CLARO  = RGBColor(0xFF, 0xD7, 0xD7)

W, H     = Inches(13.33), Inches(7.5)
DATA_DIR = "/tmp/raul_new/todo/Raul_23_05_2026"

def nueva_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def fondo(slide, color=RGBColor(0xF9,0xF9,0xF9)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_txt(tf, texto, bold=False, size=12, color=GRIS_OSCURO, align=PP_ALIGN.LEFT, clear=True):
    p = tf.paragraphs[0] if clear else tf.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = str(texto)
    run.font.bold  = bold
    run.font.size  = Pt(size)
    run.font.color.rgb = color

def add_txt(tf, texto, bold=False, size=12, color=GRIS_OSCURO, align=PP_ALIGN.LEFT):
    set_txt(tf, texto, bold, size, color, align, clear=False)

def banda_titulo(slide, titulo, subtitulo=None):
    banda = slide.shapes.add_shape(1, 0, 0, W, Inches(1.3))
    banda.fill.solid(); banda.fill.fore_color.rgb = AZUL
    banda.line.fill.background()
    tf = banda.text_frame; tf.word_wrap = True
    tf.margin_top = Pt(6); tf.margin_left = Inches(0.3)
    set_txt(tf, titulo, bold=True, size=24, color=BLANCO)
    if subtitulo:
        add_txt(tf, subtitulo, size=13, color=RGBColor(0xBD,0xD7,0xEE))

def celda(tabla, row, col, texto, bold=False, size=10,
          bg=None, fg=GRIS_OSCURO, align=PP_ALIGN.CENTER):
    cell = tabla.cell(row, col)
    if bg:
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = str(texto)
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = fg

# ── Cargar datos ──────────────────────────────────────────────────────────────
def cargar(fname):
    return json.load(open(os.path.join(DATA_DIR, fname)))

d_m4    = cargar("resultado_comparativa_m4.json")
d_m13   = cargar("resultado_comparativa_m13.json")
d_m14   = cargar("resultado_comparativa_m14.json")
d_m10h  = cargar("resultado_m10_hpo_kgcp.json")

ALGO_META = {
    "M4":      {"algo": "SMAC BlackBox (GP+EI)",  "familia": "Bayesiana — GP",  "color": RGBColor(0xD9,0xEA,0xF7)},
    "M10-HPO": {"algo": "SK-KGCP (100 eval)",     "familia": "Bayesiana — SK",  "color": RGBColor(0xC5,0xDE,0xF2)},
    "M13":     {"algo": "SPSA",                   "familia": "Gradiente Estoc.","color": RGBColor(0xE2,0xF0,0xD9)},
    "M14":     {"algo": "ALOE",                   "familia": "Gradiente Estoc.","color": RGBColor(0xFF,0xEB,0xCC)},
}

DATOS = {
    "M4":      d_m4,
    "M10-HPO": d_m10h,
    "M13":     d_m13,
    "M14":     d_m14,
}

PARAMS_ALGO = {
    "M4":      "n_trials=20, seed=42\nSurrogado: Gaussian Process (Matérn 5/2)\nAdquisición: Expected Improvement (EI)",
    "M10-HPO": "n_trials=100, seed=42\nSurrogado: SK + Knowledge Gradient CP\nAdquisición: KG look-ahead (HPO: hyperparáms optimizados)",
    "M13":     "max_iter=30, seed=42\nGradiente: estimación simultánea (SPSA, 2 eval/iter)\nPaso: α_k=a/(A+k)^0.602, ck=c/k^0.101",
    "M14":     "max_iter=30, seed=42\nGradiente: diferencias finitas (2d eval/iter)\nBúsqueda de línea: Armijo con reducción exponencial",
}

INCUMBENTE_LABELS = [
    ("cupos_ecografia_matrona",    "cupos_eco_matrona"),
    ("cupos_ecografia_ugd",        "cupos_eco_ugd"),
    ("cupos_laboratorio_ugd",      "cupos_lab_ugd"),
    ("dias_publicacion",           "dias_publicacion"),
    ("horas_control_post",         "horas_control_post"),
    ("horas_especialista_1ra",     "horas_especialista_1ra"),
    ("num_agentes_ugd",            "num_agentes_ugd"),
    ("num_matronas",               "num_matronas"),
    ("pct_bloqueo_1ra",            "pct_bloqueo_1ra"),
    ("pct_bloqueo_post_control",   "pct_bloqueo_post"),
    ("pct_consultas_vacias",       "pct_consultas_vacias"),
    ("pct_no_contactabilidad",     "pct_no_contactabilidad"),
]

MODULOS = ["M4", "M10-HPO", "M13", "M14"]

prs   = nueva_prs()
blank = prs.slide_layouts[6]

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s, AZUL)
bx = s.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.3), Inches(2.6))
set_txt(bx.text_frame,
        "Comparativa de Algoritmos de Optimización\nde Caja Negra — PC Raúl",
        bold=True, size=36, color=BLANCO, align=PP_ALIGN.CENTER)
bx2 = s.shapes.add_textbox(Inches(1), Inches(4.1), Inches(11.3), Inches(0.8))
set_txt(bx2.text_frame,
        "Simulador DES Clínica de Ginecología — Módulos M4, M10-HPO, M13, M14",
        size=20, color=RGBColor(0xBD,0xD7,0xEE), align=PP_ALIGN.CENTER)
bx3 = s.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11.3), Inches(0.7))
set_txt(bx3.text_frame,
        "Mayo 2026  ·  Raúl Araneda  ·  Entorno: PC local  ·  seed=42",
        size=13, color=RGBColor(0x9D,0xC3,0xE6), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — DESCRIPCIÓN DEL PROBLEMA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Descripción del Problema",
             "Optimización caja negra sobre simulador DES de Clínica de Ginecología")

bloques = [
    ("Objetivo",
     "Minimizar tts_full_days_mean: tiempo promedio total de espera (días)\ndesde derivación hasta alta en Clínica de Ginecología."),
    ("Variables de decisión (12)",
     "cupos_eco_matrona, cupos_eco_ugd, cupos_lab_ugd, dias_publicacion,\nhoras_control_post, horas_especialista_1ra, num_agentes_ugd, num_matronas,\npct_bloqueo_1ra, pct_bloqueo_post_control, pct_consultas_vacias, pct_no_contactabilidad"),
    ("Simulador",
     "simulador_clinica_baseline.py — modelo de eventos discretos (DES) estocástico con simpy.\nCada evaluación ejecuta 2 réplicas con semillas aleatorias → función ruidosa σ≈20–40 días."),
    ("Línea base vs mejor encontrado (datos Raúl)",
     "Configuración real: ~270 días  →  M10-HPO (SK-KGCP, 100 eval): 212.03 días  (−21.7%)\n"
     "M4 (GP+EI, 20 eval): 230.80 días  |  M13 (SPSA, 30 iter): 219.38 días\n"
     "M14 (ALOE, 30 iter): 270.73 días — sin mejora (búsqueda de línea Armijo rechazada siempre)"),
]
y = 1.44
for tit, cuerpo in bloques:
    bar = s.shapes.add_shape(1, Inches(0.38), Inches(y+0.08), Inches(0.07), Inches(0.72))
    bar.fill.solid(); bar.fill.fore_color.rgb = AZUL; bar.line.fill.background()
    bx = s.shapes.add_textbox(Inches(0.58), Inches(y), Inches(12.4), Inches(1.15))
    tf = bx.text_frame; tf.word_wrap = True
    set_txt(tf, tit, bold=True, size=12, color=AZUL)
    add_txt(tf, cuerpo, size=10.5, color=GRIS_OSCURO)
    y += 1.42

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CLASIFICACIÓN Y REFERENCIAS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Algoritmos Evaluados — Clasificación y Referencias")

cols = ["Módulo", "Algoritmo", "Familia", "Referencia clave"]
filas = [
    ["M4",       "SMAC BlackBox (GP+EI)",   "Bayesiana — GP",   "Hutter et al. 2011 (JAIR)"],
    ["M10-HPO",  "SK-KGCP (100 eval)",      "Bayesiana — SK",   "Scott, Powell & Frazier 2011 (SIAM J. Optim.)"],
    ["M13",      "SPSA",                    "Grad. Estocástico", "Spall 1992 (IEEE Trans. Autom. Control)"],
    ["M14",      "ALOE",                    "Grad. Estocástico", "Lim et al. 2012 (European J. Oper. Research)"],
]
BG = [RGBColor(0xD9,0xEA,0xF7), RGBColor(0xC5,0xDE,0xF2),
      RGBColor(0xE2,0xF0,0xD9), RGBColor(0xFF,0xEB,0xCC)]

t = s.shapes.add_table(5, 4, Inches(0.5), Inches(1.48), Inches(12.33), Inches(3.6)).table
for i, w in enumerate([0.9, 2.5, 1.7, 6.23]):
    t.columns[i].width = Inches(w)
for c, h in enumerate(cols):
    celda(t, 0, c, h, bold=True, size=12, bg=AZUL, fg=BLANCO)
for r, fila in enumerate(filas, 1):
    for c, val in enumerate(fila):
        aln = PP_ALIGN.LEFT if c in (1, 3) else PP_ALIGN.CENTER
        celda(t, r, c, val, size=12, bg=BG[r-1], align=aln)

y_desc = 5.25
desc_bloques = [
    ("Familia Bayesiana — GP:",
     "Proceso Gaussiano como surrogado. Predice costo en zonas no exploradas y elige el punto con mayor ganancia esperada (EI)."),
    ("Familia Bayesiana — SK:",
     "Stochastic Kriging: modela explícitamente la varianza del ruido del simulador. Con más evaluaciones supera al GP."),
    ("Familia Grad. Estocástico:",
     "Estima gradientes usando evaluaciones del simulador y sigue su descenso. Paralizable pero sensible al ruido del simulador."),
]
for tit_d, cuerpo_d in desc_bloques:
    bx_d = s.shapes.add_textbox(Inches(0.5), Inches(y_desc), Inches(12.33), Inches(0.45))
    tf_d = bx_d.text_frame; tf_d.word_wrap = True
    set_txt(tf_d, tit_d, bold=True, size=10, color=AZUL)
    add_txt(tf_d, cuerpo_d, size=9.5, color=GRIS_OSCURO)
    y_desc += 0.52

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PARÁMETROS DE CONFIGURACIÓN
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Parámetros de Configuración por Módulo",
             "Entorno: PC Raúl — seed=42 — simulador DES con 2 réplicas por evaluación")

positions_4 = [(0.28, 1.48), (4.76, 1.48), (0.28, 4.28), (4.76, 4.28)]
papers_4 = {
    "M4":      "Hutter et al. 2011 (JAIR)",
    "M10-HPO": "Scott, Powell & Frazier 2011 (SIAM)",
    "M13":     "Spall 1992 (IEEE Trans. Autom. Control)",
    "M14":     "Lim et al. 2012 (European J. Oper. Research)",
}
for idx, m in enumerate(MODULOS):
    cx, cy = positions_4[idx]
    meta = ALGO_META[m]
    rect = s.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(4.3), Inches(2.6))
    rect.fill.solid(); rect.fill.fore_color.rgb = meta["color"]
    rect.line.color.rgb = AZUL
    bx = s.shapes.add_textbox(Inches(cx+0.12), Inches(cy+0.1), Inches(4.06), Inches(2.38))
    tf = bx.text_frame; tf.word_wrap = True
    set_txt(tf, f"{m} — {meta['algo']}", bold=True, size=11, color=AZUL)
    add_txt(tf, PARAMS_ALGO[m], size=10, color=GRIS_OSCURO)
    add_txt(tf, f"\nRef: {papers_4[m]}", size=8.5, color=GRIS_MEDIO)

# Right column: tiempo y costo resumen
for idx, m in enumerate(MODULOS):
    cx, cy = 9.45, 1.48 + idx * 1.49
    d = DATOS[m]
    costo = d['costo_incumbente']
    tiempo = d['tiempo_seg'] / 3600
    rect2 = s.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(3.6), Inches(1.35))
    rect2.fill.solid(); rect2.fill.fore_color.rgb = ALGO_META[m]["color"]
    rect2.line.color.rgb = AZUL
    bx2 = s.shapes.add_textbox(Inches(cx+0.1), Inches(cy+0.08), Inches(3.38), Inches(1.2))
    tf2 = bx2.text_frame; tf2.word_wrap = True
    set_txt(tf2, m, bold=True, size=11, color=AZUL)
    add_txt(tf2, f"Costo: {costo:.2f} días", bold=True, size=13, color=VERDE if costo < 220 else (ROJO if costo > 265 else GRIS_OSCURO))
    add_txt(tf2, f"Tiempo: {tiempo:.2f} h", size=10, color=GRIS_OSCURO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TABLA RESULTADOS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Resultados — Mejor Costo Encontrado por Módulo",
             "Línea base: ~270 días  ·  Entorno: PC Raúl  ·  seed=42")

cols_r = ["Rank", "Módulo", "Algoritmo", "Familia", "Costo\n(días)", "vs Línea Base",
          "Tiempo\n(h)", "Eval.\n(n)", "Converge"]

orden = sorted(MODULOS, key=lambda m: DATOS[m]['costo_incumbente'])
converge_info = {
    "M4":      "Sí (GP)",
    "M10-HPO": "Sí (SK)",
    "M13":     "Sí (parcial)",
    "M14":     "No (Armijo=False)",
}

t_r = s.shapes.add_table(len(MODULOS)+1, 9, Inches(0.25), Inches(1.48), Inches(12.83), Inches(4.5)).table
for i, w in enumerate([0.5, 0.85, 2.1, 1.4, 1.0, 1.1, 0.85, 0.75, 1.5]):
    t_r.columns[i].width = Inches(w)
for c, h in enumerate(cols_r):
    celda(t_r, 0, c, h, bold=True, size=10.5, bg=AZUL, fg=BLANCO)

for r, m in enumerate(orden, 1):
    d     = DATOS[m]
    bg    = ALGO_META[m]["color"]
    costo = d['costo_incumbente']
    mejora = f"−{round((270.73-costo)/270.73*100,1)}%"
    tiempo = round(d['tiempo_seg']/3600, 2)
    n_eval = d.get('n_evaluaciones') or (d.get('max_iter', 30) * 2 if m in ("M13","M14") else "—")
    conv   = converge_info[m]
    vals   = [str(r), m, ALGO_META[m]["algo"], ALGO_META[m]["familia"],
              f"{costo:.2f}", mejora, str(tiempo), str(n_eval), conv]
    for c, val in enumerate(vals):
        aln = PP_ALIGN.LEFT if c in (2, 3, 8) else PP_ALIGN.CENTER
        fg  = VERDE if (c == 4 and r == 1) else (ROJO if (c == 4 and costo >= 265) else GRIS_OSCURO)
        celda(t_r, r, c, val, size=10.5, bg=bg, align=aln, fg=fg)

nota_r = s.shapes.add_textbox(Inches(0.25), Inches(6.1), Inches(12.83), Inches(0.65))
tf_nota = nota_r.text_frame; tf_nota.word_wrap = True
set_txt(tf_nota,
        "★ Mejor resultado: M10-HPO SK-KGCP con 100 evaluaciones → 212.03 días (−21.7% vs línea base).",
        bold=True, size=10.5, color=VERDE)
add_txt(tf_nota,
        "M14 ALOE no convergió: búsqueda de línea Armijo rechazada en todos los 30 pasos → quedó en x0=270.73 días.",
        size=10, color=ROJO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TABLA INCUMBENTES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Parámetros Óptimos (Incumbentes) por Módulo",
             "Configuración que minimizó tts_full_days_mean en cada algoritmo")

nrows_i = len(INCUMBENTE_LABELS) + 1
ncols_i = len(MODULOS) + 1
t_i = s.shapes.add_table(nrows_i, ncols_i, Inches(0.15), Inches(1.45), Inches(13.03), Inches(5.75)).table
t_i.columns[0].width = Inches(2.4)
for i in range(len(MODULOS)):
    t_i.columns[i+1].width = Inches(2.65)

celda(t_i, 0, 0, "Parámetro", bold=True, size=11, bg=AZUL, fg=BLANCO, align=PP_ALIGN.LEFT)
for c, m in enumerate(MODULOS, 1):
    costo = DATOS[m]['costo_incumbente']
    celda(t_i, 0, c, f"{m}\n({costo:.2f} días)", bold=True, size=11, bg=AZUL, fg=BLANCO)

for r, (key, label) in enumerate(INCUMBENTE_LABELS, 1):
    bg_row = AZUL_CLARO if r % 2 == 0 else BLANCO
    celda(t_i, r, 0, label, size=10, bg=bg_row, align=PP_ALIGN.LEFT)
    vals = [DATOS[m]['incumbente'].get(key, '—') for m in MODULOS]
    float_vals = [v for v in vals if isinstance(v, float)]
    for c, (m, v) in enumerate(zip(MODULOS, vals), 1):
        val_str = f"{v:.4f}" if isinstance(v, float) else str(v)
        is_min  = isinstance(v, float) and float_vals and v == min(float_vals)
        fg = VERDE if (("pct" in key) and is_min and m != "M14") else GRIS_OSCURO
        celda(t_i, r, c, val_str, size=10, bg=bg_row, fg=fg)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ANÁLISIS DE INCUMBENTES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Análisis de Incumbentes — Patrones y Diferencias",
             "¿Qué variables mueven los algoritmos? ¿Por qué difieren?")

patrones = [
    ("horas_especialista_1ra: M4=19h, M10-HPO=22h, M13=30h (máx), M14=19h (sin cambio)",
     "Los bayesianos usan horas moderadas — el SPSA (M13) explora más agresivamente y llega al máximo.\nM14 (ALOE) no se mueve de la configuración inicial x0."),
    ("pct_bloqueo_1ra: M4=8.9%, M10-HPO=15.1%, M13=5% (mín), M14=27.5% (inicial)",
     "Todos los algoritmos que convergen reducen bloqueos. M13 SPSA llega al mínimo absoluto (5%).\nM14 no optimizó nada — mantiene el punto inicial donde el ruido domina."),
    ("cupos_laboratorio_ugd: M4=67, M10-HPO=81, M13=100 (máx), M14=60",
     "SPSA satura el laboratorio al máximo. KGCP usa 81. M4 usa 67. M14 sin cambio."),
    ("Diferencia M4 (230.8d) vs M10-HPO (212.0d): +80 evaluaciones valen la pena",
     "Con 100 eval el SK-KGCP supera significativamente al GP básico con 20 eval.\nEl surrogado más sofisticado necesita más datos para mostrar su ventaja."),
    ("M14 ALOE: Armijo rechazado en 30/30 iteraciones → x_final = x0",
     "El ruido σ≈20-40 días supera la señal del gradiente estimado.\nCada paso parece 'ascender' pero es artefacto del ruido → Armijo dice siempre False."),
]

y_p = 1.44
for tit_p, cuerpo_p in patrones:
    bar = s.shapes.add_shape(1, Inches(0.35), Inches(y_p+0.08), Inches(0.07), Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = AZUL; bar.line.fill.background()
    bx = s.shapes.add_textbox(Inches(0.55), Inches(y_p), Inches(12.43), Inches(0.93))
    tf = bx.text_frame; tf.word_wrap = True
    set_txt(tf, tit_p, bold=True, size=10.5, color=AZUL)
    add_txt(tf, cuerpo_p, size=10, color=GRIS_OSCURO)
    y_p += 1.1

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — GRÁFICAS COMPARATIVAS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Gráficas Comparativas — Datos PC Raúl",
             "Convergencia, configuraciones incumbentes y tabla resumen generadas por comparativa")

comp_pngs = [
    ("Convergencia", os.path.join(DATA_DIR, "comparativa_convergencia.png")),
    ("Configuraciones", os.path.join(DATA_DIR, "comparativa_configuraciones.png")),
    ("Tabla Resumen", os.path.join(DATA_DIR, "comparativa_tabla_resumen.png")),
]
positions_c = [(0.2, 1.38, 4.25, 5.9), (4.53, 1.38, 4.25, 5.9), (8.86, 1.38, 4.25, 5.9)]
for (lbl, path), (x, y, w, h) in zip(comp_pngs, positions_c):
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        bx = s.shapes.add_textbox(Inches(x), Inches(y+2.5), Inches(w), Inches(0.5))
        set_txt(bx.text_frame, f"[{lbl}: no disponible]", size=10, color=NARANJA, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — GRÁFICAS INDIVIDUALES M13, M14, M10-HPO
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Gráficas de Convergencia Individuales — M13, M14, M10-HPO",
             "Evolución del incumbente por iteración/evaluación")

ind_pngs = [
    ("M13 SPSA",      os.path.join(DATA_DIR, "resultado_comparativa_m13.png")),
    ("M14 ALOE",      os.path.join(DATA_DIR, "resultado_comparativa_m14.png")),
    ("M10-HPO KGCP",  os.path.join(DATA_DIR, "resultado_m10_hpo_kgcp.png")),
]
for (lbl, path), (x, y, w, h) in zip(ind_pngs, positions_c):
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        bx = s.shapes.add_textbox(Inches(x), Inches(y+2.5), Inches(w), Inches(0.5))
        set_txt(bx.text_frame, f"[{lbl}: no disponible]", size=10, color=NARANJA, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — DIAGNÓSTICO M14 ALOE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s, ROJO_CLARO)
banda_titulo(s, "M14 ALOE — Diagnóstico: Búsqueda de Línea Paralizada",
             "Por qué ALOE no mejoró la solución inicial en 30 iteraciones")

diag_bloques = [
    ("Síntoma observado",
     "costo_incumbente = 270.73 días = x0 inicial. Las 30 iteraciones devuelven Armijo=False.\n"
     "El algoritmo nunca actualiza la solución — queda paralizado en el punto de partida."),
    ("Criterio de Armijo (búsqueda de línea)",
     "Armijo acepta el paso α_k·d_k si: f(x + α_k·d_k) < f(x) − σ·α_k·||∇f||²\n"
     "Cuando el simulador es ruidoso, f(x + α·d) puede ser MAYOR que f(x) por puro azar,\n"
     "incluso si la dirección de descenso d_k es correcta."),
    ("Causa raíz: relación señal/ruido",
     "Ruido del simulador: σ ≈ 20–40 días (estimado de comparativa M13/M4).\n"
     "Gradiente estimado con diff. finitas: ||∇f|| ≈ 10–20 días/unidad → señal < ruido.\n"
     "α_k se reduce exponencialmente: 1.0 → 0.8 → ... → 0.0012 → todas rechazadas."),
    ("Convergencia M13 SPSA vs M14 ALOE",
     "SPSA (M13): usa perturbaciones aleatorias SIMULTÁNEAS (Bernoulli) → promedia el ruido.\n"
     "              Llega a 219.4 días en 30 iter (8.65 h).\n"
     "ALOE (M14): usa diferencias finitas COORDENADAS → 2d=24 eval/iter → más caro, más ruido.\n"
     "              No mejora. 69.5 h gastadas sin converger."),
    ("Conclusión",
     "ALOE requiere nivel de ruido bajo para que la búsqueda de línea funcione.\n"
     "Para simuladores DES con alta varianza: preferir SPSA (M13) o métodos Bayesianos (M4, M10-HPO)."),
]

y_d = 1.4
for tit_d, cuerpo_d in diag_bloques:
    bar = s.shapes.add_shape(1, Inches(0.35), Inches(y_d+0.06), Inches(0.07), Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = ROJO; bar.line.fill.background()
    bx = s.shapes.add_textbox(Inches(0.55), Inches(y_d), Inches(12.43), Inches(0.88))
    tf = bx.text_frame; tf.word_wrap = True
    set_txt(tf, tit_d, bold=True, size=10.5, color=ROJO)
    add_txt(tf, cuerpo_d, size=10, color=GRIS_OSCURO)
    y_d += 1.0

# ══════════════════════════════════════════════════════════════════════════════
# SLIDES 11–14 — RESUMEN "PARA DUMMIES" DE CADA MÓDULO
# ══════════════════════════════════════════════════════════════════════════════
DUMMIES = [
    {
        "modulo":    "M4",
        "algo":      "SMAC BlackBox — Gaussian Process + EI",
        "paper":     "Hutter, Hoos & Leyton-Brown (2011) — Sequential Model-Based Optimization for General Algorithm Configuration — Journal of Artificial Intelligence Research (JAIR)",
        "analogia":  "Un sumiller que aprende tus gustos probando pocas botellas",
        "idea_simple": (
            "Quieres encontrar la mejor configuración de la clínica, pero simular es caro.\n"
            "SMAC prueba 20 configuraciones y construye un 'mapa mental' (modelo GP) que predice\n"
            "cuántos días de espera generaría cualquier configuración no probada aún.\n\n"
            "En cada paso elige la configuración donde la ganancia esperada (Expected Improvement)\n"
            "es mayor: balance entre explorar zonas desconocidas y explotar zonas buenas conocidas.\n\n"
            "El GP es como una superficie curva en 12 dimensiones que se va afinando\n"
            "con cada nueva evaluación del simulador."
        ),
        "cuando_usar": "≤50 evaluaciones disponibles, función ruidosa, muchas variables (hasta ~20).",
        "fortaleza": "Máxima eficiencia: mejor resultado con mínimas evaluaciones. Robusto al ruido.",
        "debilidad": "El GP escala mal a >20 variables. No modela heterocedasticidad (ruido variable en espacio).",
        "resultado": "230.80 días (PC Raúl) con 20 eval y 3.09 h. [Nube: 205.17 días — diferencia por semilla/hardware]",
        "color_bg":  RGBColor(0xD9,0xEA,0xF7),
    },
    {
        "modulo":    "M10-HPO",
        "algo":      "SK-KGCP — Knowledge Gradient with Correlated Prior (100 eval, HPO)",
        "paper":     "Scott, Powell & Frazier (2011) — The Correlated Knowledge Gradient for Simulation Optimization of Continuous Parameters Using Gaussian Process Regression — SIAM Journal on Optimization",
        "analogia":  "El sumiller que, antes de pedir una botella, piensa en cómo cambiaría su ranking completo",
        "idea_simple": (
            "Knowledge Gradient (KG) es la política óptima en problemas de 1 paso adelante:\n"
            "elige el punto que maximiza el valor esperado del MEJOR resultado final\n"
            "después de hacer esa evaluación, considerando toda la correlación del espacio.\n\n"
            "Es como ajedrez: no solo busca la mejor jugada ahora, sino la que deja\n"
            "la mejor posición para la siguiente. El surrogado es Stochastic Kriging (SK)\n"
            "que además modela la varianza del ruido del simulador.\n\n"
            "HPO: los hiperparámetros del surrogado se optimizaron con Bayesian HPO,\n"
            "dando una calibración más precisa del modelo."
        ),
        "cuando_usar": "Presupuesto alto (50–200 eval). Función suave con correlación espacial fuerte.",
        "fortaleza": "Teóricamente óptimo en 1 paso look-ahead. Con 100 eval es el mejor resultado: 212.03 días.",
        "debilidad": "Con solo 20 eval no muestra su ventaja. Costoso computacionalmente (cálculo de KG).",
        "resultado": "212.03 días — MEJOR resultado global (PC Raúl) con 100 eval y 12.27 h.",
        "color_bg":  RGBColor(0xC5,0xDE,0xF2),
    },
    {
        "modulo":    "M13",
        "algo":      "SPSA — Simultaneous Perturbation Stochastic Approximation",
        "paper":     "Spall, J.C. (1992) — Multivariate Stochastic Approximation Using a Simultaneous Perturbation Gradient Approximation — IEEE Transactions on Automatic Control, 37(3), 332–341",
        "analogia":  "El explorador que sacude todo el mapa a la vez en lugar de moverse un paso a la vez",
        "idea_simple": (
            "SPSA estima el gradiente en 12 dimensiones usando SOLO 2 evaluaciones del simulador\n"
            "por iteración — sin importar cuántas variables haya.\n\n"
            "Genera un vector aleatorio Δ de ±1 (distribución de Bernoulli) y evalúa:\n"
            "   f(x + c·Δ)  y  f(x − c·Δ)\n"
            "La diferencia dividida por 2c·Δ estima el gradiente simultáneamente en todas las dimensiones.\n\n"
            "Como el gradiente tiene ruido, el paso α_k decrece lentamente según k.\n"
            "Converge a un mínimo local con alta probabilidad si el ruido es aditivo."
        ),
        "cuando_usar": "Muchas variables (d>10), presupuesto evaluaciones limitado, función con ruido moderado.",
        "fortaleza": "Solo 2 eval/iter sin importar d. Robusto a ruido moderado. Convergencia probada.",
        "debilidad": "Lento: 30 iter × 2 eval = solo 60 eval netas, pero toma 8.65 h por simulaciones lentas.",
        "resultado": "219.38 días en 30 iter (8.65 h). Mejor que M4 en días pero mucho más tiempo.",
        "color_bg":  RGBColor(0xE2,0xF0,0xD9),
    },
    {
        "modulo":    "M14",
        "algo":      "ALOE — Armijo Line-search Optimization with Estimates",
        "paper":     "Lim, Cao & Shi (2012) — Simulation Optimization with Noisy Function Evaluations — European Journal of Operational Research, 220(3), 684–697",
        "analogia":  "El excursionista con brújula que comprueba dos veces cada paso antes de moverse, pero la niebla lo confunde",
        "idea_simple": (
            "ALOE estima el gradiente con diferencias finitas coordenadas (2d evals/iter)\n"
            "y luego hace una búsqueda de línea tipo Armijo para elegir el tamaño del paso.\n\n"
            "Armijo comprueba: ¿el paso α me da una mejora real?\n"
            "Si f(x + α·d) < f(x) − σ·α·||∇f||²  → acepta el paso.\n"
            "Si no → reduce α por factor r<1 y vuelve a intentar.\n\n"
            "El problema: cuando el simulador es muy ruidoso, f(x+α·d) puede parecer peor\n"
            "que f(x) por puro azar aunque d apunte cuesta abajo → Armijo rechaza → α→0 → parado."
        ),
        "cuando_usar": "Funciones de bajo ruido con gradiente computado exacto o casi exacto.",
        "fortaleza": "Con ruido bajo, la búsqueda de línea garantiza descenso monotónico y convergencia.",
        "debilidad": "Con ruido alto (σ≥20 días), Armijo falla en 100% de los pasos → algoritmo paralizado.",
        "resultado": "270.73 días — SIN MEJORA. 69.5 h gastadas. x_final = x0 (Armijo rechazado 30/30 iter).",
        "color_bg":  RGBColor(0xFF,0xEB,0xCC),
    },
]

for dummy in DUMMIES:
    s = prs.slides.add_slide(blank)
    fondo(s, dummy["color_bg"])
    banda_titulo(s, f"{dummy['modulo']} — {dummy['algo']}",
                 f"Ref: {dummy['paper'][:95]}{'…' if len(dummy['paper'])>95 else ''}")

    bx_an = s.shapes.add_textbox(Inches(0.3), Inches(1.38), Inches(12.73), Inches(0.48))
    set_txt(bx_an.text_frame, f"Analogía: {dummy['analogia']}", bold=True, size=13, color=AZUL, align=PP_ALIGN.CENTER)

    bx_id = s.shapes.add_textbox(Inches(0.3), Inches(1.9), Inches(7.9), Inches(3.3))
    tf_id = bx_id.text_frame; tf_id.word_wrap = True
    set_txt(tf_id, "¿Cómo funciona? (versión simple)", bold=True, size=11, color=AZUL)
    add_txt(tf_id, dummy["idea_simple"], size=10.5, color=GRIS_OSCURO)

    div = s.shapes.add_shape(1, Inches(8.25), Inches(1.9), Inches(0.04), Inches(5.3))
    div.fill.solid(); div.fill.fore_color.rgb = AZUL; div.line.fill.background()

    bx_rt = s.shapes.add_textbox(Inches(8.4), Inches(1.9), Inches(4.65), Inches(5.3))
    tf_rt = bx_rt.text_frame; tf_rt.word_wrap = True
    first = True
    for label, valor, col in [
        ("¿Cuándo usar?",       dummy["cuando_usar"],  AZUL),
        ("Fortaleza",           dummy["fortaleza"],    VERDE),
        ("Debilidad",           dummy["debilidad"],    NARANJA),
        ("Resultado obtenido",  dummy["resultado"],    GRIS_OSCURO),
    ]:
        if first:
            set_txt(tf_rt, label, bold=True, size=10.5, color=col)
            first = False
        else:
            add_txt(tf_rt, label, bold=True, size=10.5, color=col)
        add_txt(tf_rt, valor, size=10, color=GRIS_OSCURO)
        add_txt(tf_rt, " ", size=5)

    bx_paper = s.shapes.add_textbox(Inches(0.3), Inches(5.25), Inches(7.9), Inches(0.65))
    set_txt(bx_paper.text_frame, f"Referencia completa: {dummy['paper']}", size=8.5, color=GRIS_MEDIO)

# ── Guardar ───────────────────────────────────────────────────────────────────
out = "/home/user/SO/comparativa_optimizadores_raul.pptx"
prs.save(out)
print(f"PPT guardado: {out}  ({len(prs.slides)} slides)")
