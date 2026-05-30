"""
Genera presentación PPT con tablas y resúmenes de comparativa de optimizadores.
"""
import os, json
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Paleta ────────────────────────────────────────────────────────────────────
AZUL        = RGBColor(0x1F, 0x49, 0x7D)   # encabezados
AZUL_CLARO  = RGBColor(0xBD, 0xD7, 0xEE)   # filas alternas
VERDE       = RGBColor(0x37, 0x86, 0x48)   # mejor valor
GRIS_OSCURO = RGBColor(0x40, 0x40, 0x40)
BLANCO      = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_BG     = RGBColor(0xF2, 0xF2, 0xF2)
NARANJA     = RGBColor(0xC5, 0x5A, 0x11)   # advertencia

W, H = Inches(13.33), Inches(7.5)          # 16:9 widescreen

def nueva_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def fondo(slide, color=RGBColor(0xF9,0xF9,0xF9)):
    from pptx.util import Pt
    from pptx.enum.dml import MSO_THEME_COLOR
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def txt(tf, texto, bold=False, size=12, color=GRIS_OSCURO, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(texto)
    run.font.bold  = bold
    run.font.size  = Pt(size)
    run.font.color.rgb = color

def add_txt(tf, texto, bold=False, size=12, color=GRIS_OSCURO, align=PP_ALIGN.LEFT):
    """Add paragraph to text frame"""
    p = tf.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = str(texto)
    run.font.bold  = bold
    run.font.size  = Pt(size)
    run.font.color.rgb = color

def titulo_slide(slide, titulo, subtitulo=None):
    # banda superior azul
    banda = slide.shapes.add_shape(1, 0, 0, W, Inches(1.3))
    banda.fill.solid(); banda.fill.fore_color.rgb = AZUL
    banda.line.fill.background()
    tf = banda.text_frame; tf.word_wrap = True
    tf.margin_top = Pt(6); tf.margin_left = Inches(0.3)
    txt(tf, titulo, bold=True, size=24, color=BLANCO)
    if subtitulo:
        add_txt(tf, subtitulo, size=14, color=RGBColor(0xBD,0xD7,0xEE))

def celda(tabla, row, col, texto, bold=False, size=10,
          bg=None, fg=GRIS_OSCURO, align=PP_ALIGN.CENTER):
    cell = tabla.cell(row, col)
    if bg:
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = str(texto)
    run.font.bold  = bold
    run.font.size  = Pt(size)
    run.font.color.rgb = fg

# ── Datos ─────────────────────────────────────────────────────────────────────
MODULES = {
    "M4":  {"algo": "SMAC BlackBox (GP+EI)",   "familia": "Bayesian",    "paper": "Hutter et al. 2011 (SMAC)", "params": "n_trials=20, seed=42, GP kernel=Matérn 5/2", "fuente": "nube"},
    "M7":  {"algo": "SMAC + SK (EI)",           "familia": "Bayesian-SK", "paper": "Ankenman et al. 2010 (SK) + Hutter 2011", "params": "n_trials=20, seed=42, SK heteroscedástico", "fuente": "nube"},
    "M8":  {"algo": "SK Adaptativo",            "familia": "Bayesian-SK", "paper": "Ankenman et al. 2010 (SK)", "params": "n_trials=20, seed=42, δ adaptativo", "fuente": "nube"},
    "M9":  {"algo": "SK-REVI",                  "familia": "Bayesian-SK", "paper": "Quan et al. 2013 (REVI)", "params": "n_trials=20, seed=42, β=1.0 REVI", "fuente": "nube"},
    "M10": {"algo": "SK-KGCP",                  "familia": "Bayesian-SK", "paper": "Scott et al. 2011 (KG)", "params": "n_trials=20, seed=42, KG look-ahead", "fuente": "nube"},
    "M11": {"algo": "ASTRO-DF",                 "familia": "Trust-Region","paper": "Poloczek et al. 2017 (ASTRO-DF)", "params": "max_iter=3, Δ0=0.30, η1=0.10, γ1=2.0, γ2=0.50", "fuente": "nube (en curso)"},
    "M12": {"algo": "STRONG",                   "familia": "Trust-Region","paper": "Chang et al. 2013 (STRONG)", "params": "max_iter=3, n0=10, n_r=10, η0=0.01, η1=0.30", "fuente": "nube (en curso)"},
    "M13": {"algo": "SPSA",                     "familia": "Gradient-Free","paper": "Spall 1992 (SPSA)", "params": "max_iter=30, α=0.602, γ=0.101, a=0.10, A=10, n_reps=5", "fuente": "Raúl"},
    "M14": {"algo": "ALOE",                     "familia": "Gradient-Free","paper": "Rojas-Gonzalez & Van Nieuwenhuyse 2020 (ALOE)", "params": "max_iter=30, r=5, θ=0.20, γ=0.80, α0=1.0, h=0.01", "fuente": "Raúl"},
}

RESULTS_NUBE = {
    "M4":  {"costo": 205.17, "tiempo_h": 1.32, "n_eval": 20, "convergio": True},
    "M7":  {"costo": 217.76, "tiempo_h": 2.55, "n_eval": 20, "convergio": True},
    "M8":  {"costo": 218.93, "tiempo_h": 2.17, "n_eval": 20, "convergio": True},
    "M9":  {"costo": 215.66, "tiempo_h": 1.65, "n_eval": 20, "convergio": True},
    "M10": {"costo": 219.96, "tiempo_h": 1.46, "n_eval": 20, "convergio": True},
    "M11": {"costo": "—",    "tiempo_h": "en curso", "n_eval": "—", "convergio": False},
    "M12": {"costo": "—",    "tiempo_h": "en curso", "n_eval": "—", "convergio": False},
    "M13": {"costo": "—",    "tiempo_h": "en curso", "n_eval": "—", "convergio": False},
    "M14": {"costo": "—",    "tiempo_h": "en curso", "n_eval": "—", "convergio": False},
}

RESULTS_RAUL = {
    "M4":  {"costo": 230.80, "tiempo_h": 3.09, "n_eval": 20, "max_iter": "—"},
    "M13": {"costo": 219.38, "tiempo_h": 8.65, "n_eval": "~300", "max_iter": 30},
    "M14": {"costo": 270.73, "tiempo_h": 69.5, "n_eval": "~420", "max_iter": 30},
}

INCUMBENTES = {
    "M4":  {"cupos_eco_mat": 25, "cupos_eco_ugd": 26, "cupos_lab": 92, "dias_pub": 9,  "h_control": 70, "h_esp_1ra": 16, "agentes": 1, "matronas": 2, "pct_bloqueo_1ra": 0.050, "pct_bloqueo_post": 0.077, "pct_vacias": 0.050, "pct_no_cont": 0.050},
    "M7":  {"cupos_eco_mat": 25, "cupos_eco_ugd": 27, "cupos_lab": 71, "dias_pub": 3,  "h_control": 51, "h_esp_1ra": 16, "agentes": 1, "matronas": 2, "pct_bloqueo_1ra": 0.143, "pct_bloqueo_post": 0.103, "pct_vacias": 0.159, "pct_no_cont": 0.059},
    "M8":  {"cupos_eco_mat": 25, "cupos_eco_ugd": 26, "cupos_lab": 70, "dias_pub": 3,  "h_control": 51, "h_esp_1ra": 16, "agentes": 1, "matronas": 2, "pct_bloqueo_1ra": 0.150, "pct_bloqueo_post": 0.105, "pct_vacias": 0.149, "pct_no_cont": 0.054},
    "M9":  {"cupos_eco_mat": 24, "cupos_eco_ugd": 27, "cupos_lab": 70, "dias_pub": 3,  "h_control": 51, "h_esp_1ra": 16, "agentes": 1, "matronas": 2, "pct_bloqueo_1ra": 0.151, "pct_bloqueo_post": 0.099, "pct_vacias": 0.149, "pct_no_cont": 0.060},
    "M10": {"cupos_eco_mat": 25, "cupos_eco_ugd": 26, "cupos_lab": 70, "dias_pub": 3,  "h_control": 51, "h_esp_1ra": 16, "agentes": 1, "matronas": 2, "pct_bloqueo_1ra": 0.150, "pct_bloqueo_post": 0.105, "pct_vacias": 0.149, "pct_no_cont": 0.054},
    "M13_raul": {"cupos_eco_mat": 50, "cupos_eco_ugd": 10, "cupos_lab": 100, "dias_pub": 1, "h_control": 20, "h_esp_1ra": 30, "agentes": 1, "matronas": 1, "pct_bloqueo_1ra": 0.050, "pct_bloqueo_post": 0.050, "pct_vacias": 0.500, "pct_no_cont": 0.050},
    "M14_raul": {"cupos_eco_mat": 30, "cupos_eco_ugd": 30, "cupos_lab": 60, "dias_pub": 6, "h_control": 45, "h_esp_1ra": 19, "agentes": 2, "matronas": 2, "pct_bloqueo_1ra": 0.275, "pct_bloqueo_post": 0.275, "pct_vacias": 0.275, "pct_no_cont": 0.275},
}

prs = nueva_prs()
blank = prs.slide_layouts[6]  # completamente en blanco

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s, AZUL)
# título principal
box = s.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.3), Inches(2))
tf  = box.text_frame; tf.word_wrap = True
txt(tf, "Comparativa de Algoritmos de Optimización\nde Caja Negra", bold=True, size=36, color=BLANCO, align=PP_ALIGN.CENTER)
# subtítulo
box2 = s.shapes.add_textbox(Inches(1), Inches(3.9), Inches(11.3), Inches(0.9))
tf2  = box2.text_frame
txt(tf2, "Simulador DES Clínica de Ginecología — Módulos M4 a M14", size=20, color=RGBColor(0xBD,0xD7,0xEE), align=PP_ALIGN.CENTER)
# fecha y autor
box3 = s.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.3), Inches(0.6))
tf3  = box3.text_frame
txt(tf3, "Mayo 2026  ·  Raúl Araneda", size=14, color=RGBColor(0x9D,0xC3,0xE6), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
titulo_slide(s, "Contenido de la Presentación")
items = [
    "1.  Descripción del problema y simulador",
    "2.  Clasificación y referencias de algoritmos (M4–M14)",
    "3.  Parámetros de configuración por módulo",
    "4.  Resultados: métricas de desempeño (costo y tiempo)",
    "5.  Parámetros óptimos encontrados (incumbentes)",
    "6.  Gráficas de convergencia individuales",
    "7.  Gráficas comparativas (resultados Raúl PC)",
    "8.  Discusión y conclusiones",
]
box = s.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.5))
tf  = box.text_frame; tf.word_wrap = True
txt(tf, items[0], size=15, color=AZUL, bold=True)
for it in items[1:]:
    add_txt(tf, it, size=15, color=GRIS_OSCURO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DESCRIPCIÓN DEL PROBLEMA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
titulo_slide(s, "Descripción del Problema",
             "Optimización caja negra sobre simulador de eventos discretos (DES)")

bloques = [
    ("Objetivo",
     "Minimizar tts_full_days_mean: tiempo promedio total de espera (días) de pacientes\ndesde derivación hasta alta en la Clínica de Ginecología."),
    ("Variables de decisión (12)",
     "cupos_eco_mat, cupos_eco_ugd, cupos_lab_ugd, dias_publicacion, horas_control_post,\nhoras_especialista_1ra, num_agentes_ugd, num_matronas,\npct_bloqueo_1ra, pct_bloqueo_post, pct_consultas_vacias, pct_no_contactabilidad"),
    ("Simulador",
     "simulador_clinica_baseline.py — modelo DES estocástico (simpy).\nCada evaluación: 2–3 réplicas con semillas aleatorias → ruido alto σ²(x)."),
    ("Línea base",
     "Configuración real de la clínica: tts_full_days_mean ≈ 270 días.\nMejora potencial con optimización: hasta ~205 días (M4, −24 %)."),
]

y = 1.4
for titulo_blk, cuerpo in bloques:
    # barra lateral azul
    bar = s.shapes.add_shape(1, Inches(0.4), Inches(y), Inches(0.08), Inches(0.8))
    bar.fill.solid(); bar.fill.fore_color.rgb = AZUL; bar.line.fill.background()
    # texto
    bx = s.shapes.add_textbox(Inches(0.6), Inches(y-0.05), Inches(12.3), Inches(1.0))
    tf = bx.text_frame; tf.word_wrap = True
    txt(tf, titulo_blk, bold=True, size=12, color=AZUL)
    add_txt(tf, cuerpo, size=11, color=GRIS_OSCURO)
    y += 1.4

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — TABLA CLASIFICACIÓN Y REFERENCIAS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
titulo_slide(s, "Clasificación de Algoritmos y Referencias",
             "Módulos M4–M14 organizados por familia metodológica")

cols = ["Módulo", "Algoritmo", "Familia", "Referencia clave"]
filas = [
    ["M4",  "SMAC BlackBox (GP+EI)",  "Bayesiana — GP",       "Hutter et al. 2011 (JAIR)"],
    ["M7",  "SMAC + SK (EI)",          "Bayesiana — SK",       "Ankenman et al. 2010 + Hutter 2011"],
    ["M8",  "SK Adaptativo",           "Bayesiana — SK",       "Ankenman et al. 2010 (Management Sci.)"],
    ["M9",  "SK-REVI",                 "Bayesiana — SK",       "Quan et al. 2013 (IIE Trans.)"],
    ["M10", "SK-KGCP",                 "Bayesiana — SK",       "Scott et al. 2011 (Winter Sim.)"],
    ["M11", "ASTRO-DF",                "Trust-Region (lineal)","Poloczek et al. 2017 (Winter Sim.)"],
    ["M12", "STRONG",                  "Trust-Region (cuadr.)","Chang et al. 2013 (ACM TOMACS)"],
    ["M13", "SPSA",                    "Gradiente estocástico","Spall 1992 (IEEE Trans. Autom. Control)"],
    ["M14", "ALOE",                    "Gradiente estocástico","Rojas-Gonzalez & Van Nieuwenhuyse 2020"],
]

ncols = len(cols)
nrows = len(filas) + 1
table = s.shapes.add_table(nrows, ncols, Inches(0.3), Inches(1.45), Inches(12.7), Inches(5.7)).table
col_w = [0.7, 2.3, 2.1, 5.2]
for i,w in enumerate(col_w):
    table.columns[i].width = Inches(w)

for c,h in enumerate(cols):
    celda(table, 0, c, h, bold=True, size=11, bg=AZUL, fg=BLANCO)

FAMILIAS_COLOR = {
    "Bayesiana — GP":       RGBColor(0xD9,0xEA,0xF7),
    "Bayesiana — SK":       RGBColor(0xC5,0xDE,0xF2),
    "Trust-Region (lineal)":RGBColor(0xE2,0xF0,0xD9),
    "Trust-Region (cuadr.)":RGBColor(0xD5,0xEB,0xCB),
    "Gradiente estocástico":RGBColor(0xFD,0xF0,0xD0),
}
for r, fila in enumerate(filas, start=1):
    bg = FAMILIAS_COLOR.get(fila[2], RGBColor(0xF9,0xF9,0xF9))
    for c, val in enumerate(fila):
        aln = PP_ALIGN.CENTER if c in (0,2) else PP_ALIGN.LEFT
        celda(table, r, c, val, size=10, bg=bg, align=aln)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PARÁMETROS DE CONFIGURACIÓN
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
titulo_slide(s, "Parámetros de Configuración por Módulo",
             "Hiperparámetros utilizados en cada algoritmo")

params_info = [
    ("M4 — SMAC BlackBox (GP+EI)",    "n_trials=20, seed=42\nSurrogado: Gaussian Process (kernel Matérn 5/2)\nAdquisición: Expected Improvement (EI)"),
    ("M7 — SMAC + SK (EI)",           "n_trials=20, seed=42\nSurrogado: Stochastic Kriging heteroscedástico\nAdquisición: EI con corrección de ruido"),
    ("M8 — SK Adaptativo",            "n_trials=20, seed=42\nSurrogado: SK con δ adaptativo por región\nAdquisición: EI adaptativo"),
    ("M9 — SK-REVI",                  "n_trials=20, seed=42, β=1.0\nSurrogado: SK + Replicated EVI (REVI)\nAdquisición: Expected Value of Information replicado"),
    ("M10 — SK-KGCP",                 "n_trials=20, seed=42\nSurrogado: SK + Knowledge Gradient Correlated Prior\nAdquisición: KG look-ahead"),
    ("M11 — ASTRO-DF",                "max_iter=3, seed=42, Δ0=0.30\nη1=0.10, γ1=2.0, γ2=0.50, κ=0.10, μ=1.0\nn réplicas adaptativas (~30 por punto)"),
    ("M12 — STRONG",                  "max_iter=3, seed=42, Δ0=0.30\nn0=10, n_r=10, η0=0.01, η1=0.30\nModelo cuadrático local con réplicas adaptativas"),
    ("M13 — SPSA",                    "max_iter=30, n_reps=5, seed=42\nα=0.602, γ=0.101, a=0.10, A=10\nGradiente estimado con 2 evaluaciones por iter"),
    ("M14 — ALOE",                    "max_iter=30, r=5, seed=42\nθ=0.20, γ=0.80, α0=1.0, αmax=10, εf=1.0, h=0.01\nGradiente simultáneo + Armijo (2d eval/iter)"),
]

# dos columnas
n = len(params_info)
half = (n+1)//2
for idx, (titulo_p, param_txt) in enumerate(params_info):
    col_x = 0.3 + (idx // half) * 6.55
    row_y = 1.5 + (idx % half) * 1.15
    bg_color = list(FAMILIAS_COLOR.values())[idx // 2 if idx < 6 else 2 + idx - 6]
    bx = s.shapes.add_textbox(Inches(col_x), Inches(row_y), Inches(6.3), Inches(1.05))
    tf = bx.text_frame; tf.word_wrap = True
    txt(tf, titulo_p, bold=True, size=10, color=AZUL)
    add_txt(tf, param_txt, size=9, color=GRIS_OSCURO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TABLA RESULTADOS (NUBE)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
titulo_slide(s, "Resultados — Entorno Nube (Servidor Cloud)",
             "n_trials=20, max_iter=3, seed=42, n_corridas=2")

cols_r = ["Módulo", "Algoritmo", "Familia", "Mejor Costo\n(días)", "Tiempo\n(h)", "Eval.\n(n)", "Estado"]
filas_r = [
    ["M4",  "SMAC BlackBox GP+EI",  "Bayesiana", "205.17 ★", "1.32", "20", "✓ Converge"],
    ["M7",  "SMAC + SK (EI)",        "Bayesiana", "217.76",   "2.55", "20", "✓ Converge"],
    ["M8",  "SK Adaptativo",         "Bayesiana", "218.93",   "2.17", "20", "✓ Converge"],
    ["M9",  "SK-REVI",               "Bayesiana", "215.66",   "1.65", "20", "✓ Converge"],
    ["M10", "SK-KGCP",               "Bayesiana", "219.96",   "1.46", "20", "✓ Converge"],
    ["M11", "ASTRO-DF",              "Trust-Reg", "—",        "en curso","—", "⏳ Corriendo"],
    ["M12", "STRONG",                "Trust-Reg", "—",        "en curso","—", "⏳ Corriendo"],
    ["M13", "SPSA",                  "Gradiente", "—",        "en curso","—", "⏳ Corriendo"],
    ["M14", "ALOE",                  "Gradiente", "—",        "en curso","—", "⏳ Corriendo"],
]

ncols_r = len(cols_r)
nrows_r = len(filas_r) + 1
table_r = s.shapes.add_table(nrows_r, ncols_r, Inches(0.3), Inches(1.45), Inches(12.7), Inches(5.7)).table
cw_r = [0.65, 2.4, 1.3, 1.4, 1.1, 0.8, 1.55]
for i,w in enumerate(cw_r):
    table_r.columns[i].width = Inches(w)

for c,h in enumerate(cols_r):
    celda(table_r, 0, c, h, bold=True, size=10, bg=AZUL, fg=BLANCO)

ROW_BG = [
    RGBColor(0xD9,0xEA,0xF7),
    RGBColor(0xC5,0xDE,0xF2),
    RGBColor(0xC5,0xDE,0xF2),
    RGBColor(0xC5,0xDE,0xF2),
    RGBColor(0xC5,0xDE,0xF2),
    RGBColor(0xE2,0xF0,0xD9),
    RGBColor(0xE2,0xF0,0xD9),
    RGBColor(0xFD,0xF0,0xD0),
    RGBColor(0xFD,0xF0,0xD0),
]
for r, fila in enumerate(filas_r, start=1):
    bg = ROW_BG[r-1]
    for c, val in enumerate(fila):
        aln = PP_ALIGN.LEFT if c == 1 else PP_ALIGN.CENTER
        fg  = VERDE if "★" in str(val) else GRIS_OSCURO
        celda(table_r, r, c, val, size=10, bg=bg, align=aln, fg=fg)

# nota pie
nota = s.shapes.add_textbox(Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.3))
txt(nota.text_frame, "★ Mejor resultado global  ·  Línea base: ~270 días  ·  Módulos M11–M14 en ejecución al momento de esta presentación", size=9, color=GRIS_OSCURO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TABLA RESULTADOS (RAÚL PC)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
titulo_slide(s, "Resultados — Máquina Local (PC Raúl)",
             "seed=42, n_corridas=3, max_iter=30 — mayor presupuesto de iteraciones")

cols_raul = ["Módulo", "Algoritmo", "Familia", "Mejor Costo\n(días)", "Tiempo (h)", "max_iter / n_trials"]
filas_raul = [
    ["M4",  "SMAC BlackBox GP+EI",  "Bayesiana", "230.80", "3.09", "n_trials=20"],
    ["M13", "SPSA",                  "Gradiente", "219.38", "8.65", "max_iter=30"],
    ["M14", "ALOE",                  "Gradiente", "270.73", "69.5", "max_iter=30"],
]

nrows_raul = len(filas_raul) + 1
ncols_raul = len(cols_raul)
table_raul = s.shapes.add_table(nrows_raul, ncols_raul, Inches(0.5), Inches(1.55), Inches(12.3), Inches(2.6)).table
cw_raul = [0.7, 2.6, 1.3, 1.5, 1.3, 2.1]
for i,w in enumerate(cw_raul):
    table_raul.columns[i].width = Inches(w)

for c,h in enumerate(cols_raul):
    celda(table_raul, 0, c, h, bold=True, size=11, bg=AZUL, fg=BLANCO)

BG_RAUL = [
    RGBColor(0xD9,0xEA,0xF7),
    RGBColor(0xFD,0xF0,0xD0),
    RGBColor(0xFD,0xF0,0xD0),
]
for r, fila in enumerate(filas_raul, start=1):
    for c, val in enumerate(fila):
        aln = PP_ALIGN.LEFT if c == 1 else PP_ALIGN.CENTER
        celda(table_raul, r, c, val, size=11, bg=BG_RAUL[r-1], align=aln)

# recuadro explicativo
y_obs = 4.4
obs_items = [
    ("¿Por qué difieren con la nube?", AZUL, True, 13),
    ("• Estocástico por diseño: el simulador DES no tiene seed fija → cada réplica genera resultados diferentes", GRIS_OSCURO, False, 11),
    ("• n_corridas distinto: nube=2 réplicas/eval vs Raúl=3 réplicas/eval → estimados de costo distintos", GRIS_OSCURO, False, 11),
    ("• max_iter distinto: nube=3 (demo) vs Raúl=30 (completo) → más exploración del espacio en PC local", GRIS_OSCURO, False, 11),
    ("• M14 (ALOE) no mejora la línea base: convergencia prematura en espacio de alta dimensión y ruido elevado", NARANJA, False, 11),
]
box_obs = s.shapes.add_textbox(Inches(0.5), Inches(y_obs), Inches(12.3), Inches(2.8))
tf_obs  = box_obs.text_frame; tf_obs.word_wrap = True
txt(tf_obs, obs_items[0][0], bold=obs_items[0][2], size=obs_items[0][3], color=obs_items[0][1])
for ob_txt, ob_color, ob_bold, ob_sz in obs_items[1:]:
    add_txt(tf_obs, ob_txt, bold=ob_bold, size=ob_sz, color=ob_color)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TABLA INCUMBENTES (SMAC M4–M10)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
titulo_slide(s, "Parámetros Óptimos Encontrados — Módulos SMAC (M4–M10)",
             "Configuración incumbente al finalizar cada algoritmo bayesiano")

param_names = [
    "cupos_eco_mat", "cupos_eco_ugd", "cupos_lab_ugd", "dias_pub",
    "h_control", "h_esp_1ra", "agentes_ugd", "matronas",
    "pct_bloqueo_1ra", "pct_bloqueo_post", "pct_vacias", "pct_no_cont"
]
modulos_inc = ["M4", "M7", "M8", "M9", "M10"]
inc_keys    = ["cupos_eco_mat","cupos_eco_ugd","cupos_lab","dias_pub","h_control","h_esp_1ra","agentes","matronas","pct_bloqueo_1ra","pct_bloqueo_post","pct_vacias","pct_no_cont"]

nrows_i = len(param_names) + 1
ncols_i = len(modulos_inc) + 1
table_i = s.shapes.add_table(nrows_i, ncols_i, Inches(0.25), Inches(1.45), Inches(12.83), Inches(5.75)).table
cw_i = [2.33] + [2.1]*len(modulos_inc)
for i,w in enumerate(cw_i):
    table_i.columns[i].width = Inches(w)

celda(table_i, 0, 0, "Parámetro", bold=True, size=10, bg=AZUL, fg=BLANCO, align=PP_ALIGN.LEFT)
for c, m in enumerate(modulos_inc, start=1):
    costo = RESULTS_NUBE[m]["costo"]
    celda(table_i, 0, c, f"{m}\n({costo} d)", bold=True, size=10, bg=AZUL, fg=BLANCO)

for r, (pname, key) in enumerate(zip(param_names, inc_keys), start=1):
    bg = AZUL_CLARO if r % 2 == 0 else BLANCO
    celda(table_i, r, 0, pname, bold=False, size=9, bg=bg, align=PP_ALIGN.LEFT)
    vals = [INCUMBENTES[m][key] for m in modulos_inc]
    # resaltar mínimo en pct_ fields (menor = mejor control) o máximo para cupos/horas
    for c, (m, v) in enumerate(zip(modulos_inc, vals), start=1):
        fg = GRIS_OSCURO
        # resaltar si es pct que converge al mínimo (eficiencia)
        if "pct" in key and isinstance(v, float) and v == min(vals):
            fg = VERDE
        val_str = f"{v:.3f}" if isinstance(v, float) else str(v)
        celda(table_i, r, c, val_str, size=9, bg=bg, fg=fg)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — TABLA INCUMBENTES (GRADIENTE RAÚL)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
titulo_slide(s, "Parámetros Óptimos — Módulos Gradiente Estocástico (Raúl PC)",
             "M13 (SPSA) y M14 (ALOE) con max_iter=30")

modulos_g = ["M13_raul", "M14_raul"]
labels_g  = ["M13 SPSA\n(219.38 d)", "M14 ALOE\n(270.73 d)"]

nrows_g = len(param_names) + 1
ncols_g = 3
table_g = s.shapes.add_table(nrows_g, ncols_g, Inches(1.5), Inches(1.55), Inches(9.8), Inches(5.65)).table
table_g.columns[0].width = Inches(3.3)
table_g.columns[1].width = Inches(3.25)
table_g.columns[2].width = Inches(3.25)

celda(table_g, 0, 0, "Parámetro", bold=True, size=12, bg=AZUL, fg=BLANCO, align=PP_ALIGN.LEFT)
for c, lbl in enumerate(labels_g, start=1):
    celda(table_g, 0, c, lbl, bold=True, size=12, bg=AZUL, fg=BLANCO)

for r, (pname, key) in enumerate(zip(param_names, inc_keys), start=1):
    bg = AZUL_CLARO if r % 2 == 0 else BLANCO
    celda(table_g, r, 0, pname, size=11, bg=bg, align=PP_ALIGN.LEFT)
    for c, m in enumerate(modulos_g, start=1):
        v = INCUMBENTES[m][key]
        val_str = f"{v:.3f}" if isinstance(v, float) else str(v)
        celda(table_g, r, c, val_str, size=11, bg=bg)

# nota
nota9 = s.shapes.add_textbox(Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.3))
txt(nota9.text_frame, "M14 ALOE no mejoró la línea base (~270 días): convergencia prematura por ruido elevado del simulador y alta dimensionalidad.", size=9, color=NARANJA)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — GRÁFICAS INDIVIDUALES (NUBE M4, M8, M9, M10)
# ══════════════════════════════════════════════════════════════════════════════
PNGS_NUBE = {
    "M4":  "/home/user/SO/resultado_comparativa_m4.png",
    "M8":  "/home/user/SO/resultado_comparativa_m8.png",
    "M9":  "/home/user/SO/resultado_comparativa_m9.png",
    "M10": "/home/user/SO/resultado_comparativa_m10.png",
}

s = prs.slides.add_slide(blank)
fondo(s)
titulo_slide(s, "Gráficas de Convergencia Individuales — Nube (M4, M8, M9, M10)")

positions = [
    (0.25, 1.45, 6.3, 2.95),
    (6.75, 1.45, 6.3, 2.95),
    (0.25, 4.45, 6.3, 2.95),
    (6.75, 4.45, 6.3, 2.95),
]
for (m, path), (x, y, w, h) in zip(PNGS_NUBE.items(), positions):
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        bx = s.shapes.add_textbox(Inches(x), Inches(y+1), Inches(w), Inches(0.5))
        txt(bx.text_frame, f"[{m}: imagen no disponible]", size=10, color=NARANJA)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — GRÁFICAS RAÚL (M13, M14 individuales)
# ══════════════════════════════════════════════════════════════════════════════
PNGS_RAUL_IND = {
    "M13 SPSA": "/tmp/raul_results/Raul_23_05_2026/resultado_comparativa_m13.png",
    "M14 ALOE": "/tmp/raul_results/Raul_23_05_2026/resultado_comparativa_m14.png",
}

s = prs.slides.add_slide(blank)
fondo(s)
titulo_slide(s, "Gráficas de Convergencia — PC Raúl (M13 SPSA, M14 ALOE)", "max_iter=30")

pos_r = [
    (0.4,  1.45, 6.1, 5.7),
    (6.85, 1.45, 6.1, 5.7),
]
for (m, path), (x, y, w, h) in zip(PNGS_RAUL_IND.items(), pos_r):
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        bx = s.shapes.add_textbox(Inches(x), Inches(y+2), Inches(w), Inches(0.5))
        txt(bx.text_frame, f"[{m}: imagen no disponible]", size=10, color=NARANJA)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — COMPARATIVA CONVERGENCIA (RAÚL)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
titulo_slide(s, "Gráfica Comparativa — Convergencia (PC Raúl)", "M4 vs M13 vs M14")

png_conv = "/tmp/raul_results/Raul_23_05_2026/comparativa_convergencia.png"
if os.path.exists(png_conv):
    s.shapes.add_picture(png_conv, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.7))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — COMPARATIVA CONFIGURACIONES (RAÚL)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
titulo_slide(s, "Gráfica Comparativa — Configuraciones Óptimas (PC Raúl)", "M4 vs M13 vs M14")

png_conf = "/tmp/raul_results/Raul_23_05_2026/comparativa_configuraciones.png"
if os.path.exists(png_conf):
    s.shapes.add_picture(png_conf, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.7))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — COMPARATIVA TABLA RESUMEN (RAÚL)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
titulo_slide(s, "Tabla Resumen Comparativa (PC Raúl)", "Resumen visual generado por modulo_comparativa_caja_negra.py")

png_tbl = "/tmp/raul_results/Raul_23_05_2026/comparativa_tabla_resumen.png"
if os.path.exists(png_tbl):
    s.shapes.add_picture(png_tbl, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.7))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — DISCUSIÓN Y CONCLUSIONES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
titulo_slide(s, "Discusión y Conclusiones")

conclusiones = [
    ("Mejor resultado global", "M4 (SMAC GP+EI): 205.17 días, reducción del 24% respecto a la línea base (~270 días).\nPresupuesto: solo 20 evaluaciones (~1.3 h)."),
    ("Métodos Bayesianos (M4–M10)", "Todos convergen con n_trials=20. M4 supera a variantes SK pese a su mayor simplicidad.\nM7–M10 muestran desempeño similar (~216–220 días), sugiriendo que el ruido del simulador\nlimita la ventaja de surrogados más sofisticados."),
    ("Trust-Region (M11 ASTRO-DF, M12 STRONG)", "Costo computacional alto: ~2 h solo para evaluar x₀ (n≈30 réplicas adaptativas).\nAdecuados para problemas con menor ruido o mayor presupuesto de simulación."),
    ("Gradiente estocástico (M13 SPSA, M14 ALOE)", "SPSA converge a ~219 días con 30 iteraciones (~8.65 h). ALOE no mejora la línea\nbase (270 días, 69.5 h): sensible al ruido alto y alta dimensión (12 variables)."),
    ("Hallazgo clave — variables de decisión", "La palanca principal no es horas_especialista_1ra sino la eficiencia de las horas:\nreducir pct_bloqueo_1ra y pct_consultas_vacias de ~30% a ~5% equivale\na duplicar la capacidad efectiva del especialista sin cambiar dotación."),
    ("Próximos pasos", "Completar M11–M14 en la nube (max_iter=3 → resultados preliminares).\nGenerar gráficas comparativas finales con los 9 módulos juntos.\nValidar incumbentes con 30 réplicas independientes."),
]

y_c = 1.45
for titulo_c, cuerpo_c in conclusiones:
    bar = s.shapes.add_shape(1, Inches(0.35), Inches(y_c+0.05), Inches(0.08), Inches(0.65))
    bar.fill.solid(); bar.fill.fore_color.rgb = AZUL; bar.line.fill.background()
    bx = s.shapes.add_textbox(Inches(0.55), Inches(y_c), Inches(12.4), Inches(1.0))
    tf = bx.text_frame; tf.word_wrap = True
    txt(tf, titulo_c, bold=True, size=11, color=AZUL)
    add_txt(tf, cuerpo_c, size=10, color=GRIS_OSCURO)
    y_c += 1.0

# ══════════════════════════════════════════════════════════════════════════════
# GUARDAR
# ══════════════════════════════════════════════════════════════════════════════
out = "/home/user/SO/comparativa_optimizadores.pptx"
prs.save(out)
print(f"PPT guardado: {out}")
