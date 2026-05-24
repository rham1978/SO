"""
PPT v2 — incluye resultados de Raúl PC para M13/M14,
marca M11/M12 como sin datos disponibles.
"""
import os, json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paleta ────────────────────────────────────────────────────────────────────
AZUL        = RGBColor(0x1F, 0x49, 0x7D)
AZUL_CLARO  = RGBColor(0xBD, 0xD7, 0xEE)
VERDE       = RGBColor(0x37, 0x86, 0x48)
ROJO        = RGBColor(0xC0, 0x00, 0x00)
GRIS_OSCURO = RGBColor(0x40, 0x40, 0x40)
BLANCO      = RGBColor(0xFF, 0xFF, 0xFF)
NARANJA     = RGBColor(0xC5, 0x5A, 0x11)
GRIS_MEDIO  = RGBColor(0x80, 0x80, 0x80)
AMARILLO_BG = RGBColor(0xFF, 0xF2, 0xCC)

W, H = Inches(13.33), Inches(7.5)

def nueva_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def fondo(slide, color=RGBColor(0xF9,0xF9,0xF9)):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_txt(tf, texto, bold=False, size=12, color=GRIS_OSCURO, align=PP_ALIGN.LEFT, clear=True):
    if clear:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = align
    run = p.add_run()
    run.text = str(texto)
    run.font.bold  = bold
    run.font.size  = Pt(size)
    run.font.color.rgb = color

def add_txt(tf, texto, bold=False, size=12, color=GRIS_OSCURO, align=PP_ALIGN.LEFT):
    set_txt(tf, texto, bold, size, color, align, clear=False)

def banda_titulo(slide, titulo, subtitulo=None):
    banda = slide.shapes.add_shape(1, 0, 0, W, Inches(1.3))
    banda.fill.solid(); banda.fill.fore_color.rgb = AZUL
    banda.line.fill.background()
    tf = banda.text_frame; tf.word_wrap = True
    tf.margin_top = Pt(6); tf.margin_left = Inches(0.3)
    set_txt(tf, titulo, bold=True, size=24, color=BLANCO)
    if subtitulo:
        add_txt(tf, subtitulo, size=13, color=RGBColor(0xBD,0xD7,0xEE))

def celda(tabla, row, col, texto, bold=False, size=10,
          bg=None, fg=GRIS_OSCURO, align=PP_ALIGN.CENTER):
    cell = tabla.cell(row, col)
    if bg:
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
    tf = cell.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = str(texto)
    run.font.bold  = bold
    run.font.size  = Pt(size)
    run.font.color.rgb = fg

# ── Datos consolidados ────────────────────────────────────────────────────────
RESULTADOS = {
    "M4":  {"algo": "SMAC BlackBox (GP+EI)",  "familia": "Bayesiana",     "costo": 205.17, "tiempo_h": 1.32,  "n_eval": 20,    "fuente": "Nube",     "estado": "✓"},
    "M7":  {"algo": "SMAC + SK (EI)",          "familia": "Bayesiana-SK",  "costo": 217.76, "tiempo_h": 2.55,  "n_eval": 20,    "fuente": "Nube",     "estado": "✓"},
    "M8":  {"algo": "SK Adaptativo",           "familia": "Bayesiana-SK",  "costo": 218.93, "tiempo_h": 2.17,  "n_eval": 20,    "fuente": "Nube",     "estado": "✓"},
    "M9":  {"algo": "SK-REVI",                 "familia": "Bayesiana-SK",  "costo": 215.66, "tiempo_h": 1.65,  "n_eval": 20,    "fuente": "Nube",     "estado": "✓"},
    "M10": {"algo": "SK-KGCP",                 "familia": "Bayesiana-SK",  "costo": 219.96, "tiempo_h": 1.46,  "n_eval": 20,    "fuente": "Nube",     "estado": "✓"},
    "M11": {"algo": "ASTRO-DF",                "familia": "Trust-Region",  "costo": None,   "tiempo_h": None,  "n_eval": None,  "fuente": "—",        "estado": "✗"},
    "M12": {"algo": "STRONG",                  "familia": "Trust-Region",  "costo": None,   "tiempo_h": None,  "n_eval": None,  "fuente": "—",        "estado": "✗"},
    "M13": {"algo": "SPSA",                    "familia": "Gradiente",     "costo": 219.38, "tiempo_h": 8.65,  "n_eval": "~300","fuente": "Raúl PC",  "estado": "✓"},
    "M14": {"algo": "ALOE",                    "familia": "Gradiente",     "costo": 270.73, "tiempo_h": 69.5,  "n_eval": "~420","fuente": "Raúl PC",  "estado": "✗ sin mejora"},
}

PARAMS_ALGO = {
    "M4":  ("SMAC BlackBox (GP+EI)",    "n_trials=20, seed=42\nSurrogado: Gaussian Process (Matérn 5/2)\nAdquisición: Expected Improvement (EI)"),
    "M7":  ("SMAC + SK (EI)",           "n_trials=20, seed=42\nSurrogado: Stochastic Kriging heteroscedástico\nAdquisición: EI con corrección de ruido"),
    "M8":  ("SK Adaptativo",            "n_trials=20, seed=42\nSurrogado: SK con δ adaptativo por región\nAdquisición: EI adaptativo"),
    "M9":  ("SK-REVI",                  "n_trials=20, seed=42, β=1.0\nSurrogado: SK + Replicated EVI (REVI)\nAdquisición: Expected Value of Information replicado"),
    "M10": ("SK-KGCP",                  "n_trials=20, seed=42\nSurrogado: SK + Knowledge Gradient CP\nAdquisición: KG look-ahead"),
    "M11": ("ASTRO-DF",                 "max_iter=3, seed=42, Δ0=0.30\nη1=0.10, γ1=2.0, γ2=0.50, κ=0.10, μ=1.0\nn réplicas adaptativas (~30 por punto)"),
    "M12": ("STRONG",                   "max_iter=3, seed=42, Δ0=0.30\nn0=10, n_r=10, η0=0.01, η1=0.30\nModelo cuadrático local con réplicas adaptativas"),
    "M13": ("SPSA",                     "max_iter=30, n_reps=5, seed=42\nα=0.602, γ=0.101, a=0.10, A=10\nGradiente estimado con 2 evaluaciones por iter"),
    "M14": ("ALOE",                     "max_iter=30, r=5, seed=42\nθ=0.20, γ=0.80, α0=1.0, αmax=10, h=0.01\nGradiente simultáneo + Armijo (2d eval/iter)"),
}

REFS = {
    "M4":  "Hutter et al. 2011 (JAIR)",
    "M7":  "Ankenman et al. 2010 + Hutter 2011",
    "M8":  "Ankenman et al. 2010 (Management Sci.)",
    "M9":  "Quan et al. 2013 (IIE Trans.)",
    "M10": "Scott et al. 2011 (Winter Sim.)",
    "M11": "Poloczek et al. 2017 (Winter Sim.)",
    "M12": "Chang et al. 2013 (ACM TOMACS)",
    "M13": "Spall 1992 (IEEE Trans. Autom. Control)",
    "M14": "Rojas-Gonzalez & Van Nieuwenhuyse 2020",
}

INCUMBENTES = {
    "M4":  {"cupos_eco_mat":25,"cupos_eco_ugd":26,"cupos_lab":92,"dias_pub":9, "h_control":70,"h_esp_1ra":16,"agentes":1,"matronas":2,"pct_bloqueo_1ra":0.050,"pct_bloqueo_post":0.077,"pct_vacias":0.050,"pct_no_cont":0.050},
    "M7":  {"cupos_eco_mat":25,"cupos_eco_ugd":27,"cupos_lab":71,"dias_pub":3, "h_control":51,"h_esp_1ra":16,"agentes":1,"matronas":2,"pct_bloqueo_1ra":0.143,"pct_bloqueo_post":0.103,"pct_vacias":0.159,"pct_no_cont":0.059},
    "M8":  {"cupos_eco_mat":25,"cupos_eco_ugd":26,"cupos_lab":70,"dias_pub":3, "h_control":51,"h_esp_1ra":16,"agentes":1,"matronas":2,"pct_bloqueo_1ra":0.150,"pct_bloqueo_post":0.105,"pct_vacias":0.149,"pct_no_cont":0.054},
    "M9":  {"cupos_eco_mat":24,"cupos_eco_ugd":27,"cupos_lab":70,"dias_pub":3, "h_control":51,"h_esp_1ra":16,"agentes":1,"matronas":2,"pct_bloqueo_1ra":0.151,"pct_bloqueo_post":0.099,"pct_vacias":0.149,"pct_no_cont":0.060},
    "M10": {"cupos_eco_mat":25,"cupos_eco_ugd":26,"cupos_lab":70,"dias_pub":3, "h_control":51,"h_esp_1ra":16,"agentes":1,"matronas":2,"pct_bloqueo_1ra":0.150,"pct_bloqueo_post":0.105,"pct_vacias":0.149,"pct_no_cont":0.054},
    "M13": {"cupos_eco_mat":50,"cupos_eco_ugd":10,"cupos_lab":100,"dias_pub":1,"h_control":20,"h_esp_1ra":30,"agentes":1,"matronas":1,"pct_bloqueo_1ra":0.050,"pct_bloqueo_post":0.050,"pct_vacias":0.500,"pct_no_cont":0.050},
    "M14": {"cupos_eco_mat":30,"cupos_eco_ugd":30,"cupos_lab":60,"dias_pub":6, "h_control":45,"h_esp_1ra":19,"agentes":2,"matronas":2,"pct_bloqueo_1ra":0.275,"pct_bloqueo_post":0.275,"pct_vacias":0.275,"pct_no_cont":0.275},
}

PARAM_NAMES = ["cupos_eco_mat","cupos_eco_ugd","cupos_lab","dias_pub",
               "h_control","h_esp_1ra","agentes","matronas",
               "pct_bloqueo_1ra","pct_bloqueo_post","pct_vacias","pct_no_cont"]

FAMILIAS_COLOR = {
    "Bayesiana":    RGBColor(0xD9,0xEA,0xF7),
    "Bayesiana-SK": RGBColor(0xC5,0xDE,0xF2),
    "Trust-Region": RGBColor(0xE2,0xF0,0xD9),
    "Gradiente":    RGBColor(0xFD,0xF0,0xD0),
}

prs = nueva_prs()
blank = prs.slide_layouts[6]

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — PORTADA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s, AZUL)
box = s.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.3), Inches(2.2))
tf  = box.text_frame; tf.word_wrap = True
set_txt(tf, "Comparativa de Algoritmos de Optimización\nde Caja Negra", bold=True, size=36, color=BLANCO, align=PP_ALIGN.CENTER)
box2 = s.shapes.add_textbox(Inches(1), Inches(4.1), Inches(11.3), Inches(0.9))
set_txt(box2.text_frame, "Simulador DES Clínica de Ginecología — Módulos M4 a M14", size=20, color=RGBColor(0xBD,0xD7,0xEE), align=PP_ALIGN.CENTER)
box3 = s.shapes.add_textbox(Inches(1), Inches(5.4), Inches(11.3), Inches(0.5))
set_txt(box3.text_frame, "Mayo 2026  ·  Raúl Araneda  ·  Resultados combinados: Nube + PC Local", size=13, color=RGBColor(0x9D,0xC3,0xE6), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — ESTADO DE RESULTADOS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Estado de Resultados por Módulo", "Datos consolidados: Nube (M4–M10) + PC Raúl (M13, M14)")

cols = ["Módulo","Algoritmo","Familia","Costo\n(días)","Tiempo","Evaluaciones","Fuente","Estado"]
filas = [
    ["M4",  "SMAC BlackBox GP+EI", "Bayesiana",    "205.17 ★","1.32 h","20",     "Nube",    "✓"],
    ["M7",  "SMAC + SK (EI)",       "Bayesiana-SK", "217.76",  "2.55 h","20",     "Nube",    "✓"],
    ["M8",  "SK Adaptativo",        "Bayesiana-SK", "218.93",  "2.17 h","20",     "Nube",    "✓"],
    ["M9",  "SK-REVI",              "Bayesiana-SK", "215.66",  "1.65 h","20",     "Nube",    "✓"],
    ["M10", "SK-KGCP",              "Bayesiana-SK", "219.96",  "1.46 h","20",     "Nube",    "✓"],
    ["M11", "ASTRO-DF",             "Trust-Region", "—",       "—",     "—",      "—",       "✗ Sin datos"],
    ["M12", "STRONG",               "Trust-Region", "—",       "—",     "—",      "—",       "✗ Sin datos"],
    ["M13", "SPSA",                 "Gradiente",    "219.38",  "8.65 h","~300",   "Raúl PC", "✓"],
    ["M14", "ALOE",                 "Gradiente",    "270.73",  "69.5 h","~420",   "Raúl PC", "✗ Sin mejora"],
]

t = s.shapes.add_table(len(filas)+1, len(cols), Inches(0.25), Inches(1.45), Inches(12.83), Inches(5.75)).table
for i,w in enumerate([0.6,2.35,1.25,1.0,0.95,1.1,0.9,1.68]):
    t.columns[i].width = Inches(w)
for c,h in enumerate(cols):
    celda(t, 0, c, h, bold=True, size=10, bg=AZUL, fg=BLANCO)

BG_ROWS = ["Bayesiana","Bayesiana-SK","Bayesiana-SK","Bayesiana-SK","Bayesiana-SK","Trust-Region","Trust-Region","Gradiente","Gradiente"]
for r,fila in enumerate(filas, 1):
    bg = FAMILIAS_COLOR.get(BG_ROWS[r-1], RGBColor(0xF9,0xF9,0xF9))
    for c,val in enumerate(fila):
        aln = PP_ALIGN.LEFT if c==1 else PP_ALIGN.CENTER
        fg = VERDE if "★" in str(val) else (ROJO if "✗" in str(val) else GRIS_OSCURO)
        celda(t, r, c, val, size=10, bg=bg, align=aln, fg=fg)

nota = s.shapes.add_textbox(Inches(0.25), Inches(7.15), Inches(12.83), Inches(0.28))
set_txt(nota.text_frame, "★ Mejor resultado global (−24% vs línea base ~270 días)  ·  M11/M12: proceso colapsó por limitación de recursos en nube  ·  M14: no mejoró línea base", size=8.5, color=GRIS_MEDIO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DESCRIPCIÓN DEL PROBLEMA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Descripción del Problema", "Optimización caja negra sobre simulador DES de Clínica de Ginecología")

bloques = [
    ("Objetivo", "Minimizar tts_full_days_mean: tiempo promedio total de espera (días) desde derivación hasta alta."),
    ("Variables de decisión (12)", "cupos_eco_matrona, cupos_eco_ugd, cupos_lab_ugd, dias_publicacion,\nhoras_control_post, horas_especialista_1ra, num_agentes_ugd, num_matronas,\npct_bloqueo_1ra, pct_bloqueo_post_control, pct_consultas_vacias, pct_no_contactabilidad"),
    ("Simulador", "simulador_clinica_baseline.py — modelo DES estocástico (simpy). Cada evaluación ejecuta\n2–3 réplicas aleatorias → función ruidosa de alta varianza σ²(x)."),
    ("Línea base vs mejor encontrado", "Configuración real: ~270 días  →  M4 (SMAC GP+EI): 205.17 días  →  mejora del 24%\nHallazgo clave: reducir pct_bloqueo_1ra y pct_vacias de ~30% a ~5% equivale a duplicar capacidad efectiva."),
]
y = 1.42
for tit, cuerpo in bloques:
    bar = s.shapes.add_shape(1, Inches(0.38), Inches(y+0.08), Inches(0.07), Inches(0.7))
    bar.fill.solid(); bar.fill.fore_color.rgb = AZUL; bar.line.fill.background()
    bx = s.shapes.add_textbox(Inches(0.58), Inches(y), Inches(12.4), Inches(1.1))
    tf = bx.text_frame; tf.word_wrap = True
    set_txt(tf, tit, bold=True, size=12, color=AZUL)
    add_txt(tf, cuerpo, size=11, color=GRIS_OSCURO)
    y += 1.45

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — CLASIFICACIÓN Y REFERENCIAS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Clasificación de Algoritmos y Referencias Bibliográficas")

cols4 = ["Módulo","Algoritmo","Familia","Referencia clave","Datos disponibles"]
filas4 = [
    ["M4",  "SMAC BlackBox (GP+EI)", "Bayesiana — GP",       "Hutter et al. 2011 (JAIR)",                          "✓ Nube"],
    ["M7",  "SMAC + SK (EI)",         "Bayesiana — SK",       "Ankenman et al. 2010 + Hutter 2011",                 "✓ Nube"],
    ["M8",  "SK Adaptativo",          "Bayesiana — SK",       "Ankenman et al. 2010 (Management Sci.)",             "✓ Nube"],
    ["M9",  "SK-REVI",                "Bayesiana — SK",       "Quan et al. 2013 (IIE Trans.)",                      "✓ Nube"],
    ["M10", "SK-KGCP",                "Bayesiana — SK",       "Scott et al. 2011 (Winter Sim.)",                    "✓ Nube"],
    ["M11", "ASTRO-DF",               "Trust-Region (lin.)",  "Poloczek et al. 2017 (Winter Sim.)",                 "✗ Pendiente"],
    ["M12", "STRONG",                 "Trust-Region (cuad.)", "Chang et al. 2013 (ACM TOMACS)",                     "✗ Pendiente"],
    ["M13", "SPSA",                   "Gradiente estocástico","Spall 1992 (IEEE Trans. Autom. Control)",            "✓ Raúl PC"],
    ["M14", "ALOE",                   "Gradiente estocástico","Rojas-Gonzalez & Van Nieuwenhuyse 2020",             "✓ Raúl PC"],
]
t4 = s.shapes.add_table(len(filas4)+1, 5, Inches(0.25), Inches(1.45), Inches(12.83), Inches(5.75)).table
for i,w in enumerate([0.65,2.2,1.75,4.3,1.43]):
    t4.columns[i].width = Inches(w)
for c,h in enumerate(cols4):
    celda(t4, 0, c, h, bold=True, size=10, bg=AZUL, fg=BLANCO)

FAM_MAP = {"Bayesiana — GP":FAMILIAS_COLOR["Bayesiana"],"Bayesiana — SK":FAMILIAS_COLOR["Bayesiana-SK"],
           "Trust-Region (lin.)":FAMILIAS_COLOR["Trust-Region"],"Trust-Region (cuad.)":FAMILIAS_COLOR["Trust-Region"],
           "Gradiente estocástico":FAMILIAS_COLOR["Gradiente"]}
for r,fila in enumerate(filas4, 1):
    bg = FAM_MAP.get(fila[2], RGBColor(0xF9,0xF9,0xF9))
    for c,val in enumerate(fila):
        aln = PP_ALIGN.LEFT if c in (1,3) else PP_ALIGN.CENTER
        fg = VERDE if "✓" in str(val) else (ROJO if "✗" in str(val) else GRIS_OSCURO)
        celda(t4, r, c, val, size=10, bg=bg, align=aln, fg=fg)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PARÁMETROS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Parámetros de Configuración por Módulo")

items = list(PARAMS_ALGO.items())
half = 5
for idx,(m,(titulo_p, param_txt)) in enumerate(items):
    col_x = 0.28 + (idx // half) * 6.55
    row_y = 1.42 + (idx % half) * 1.17
    bx = s.shapes.add_textbox(Inches(col_x), Inches(row_y), Inches(6.3), Inches(1.1))
    tf = bx.text_frame; tf.word_wrap = True
    m_data = RESULTADOS[m]
    status_color = GRIS_OSCURO if m_data["estado"] == "✓" else NARANJA
    set_txt(tf, f"{m} — {titulo_p}", bold=True, size=10, color=AZUL)
    add_txt(tf, param_txt, size=9, color=status_color)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TABLA COMPARATIVA COMPLETA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Tabla Comparativa — Resultados Consolidados",
             "Mejor costo encontrado por cada algoritmo (línea base: ~270 días)")

# ranking visual: ordenar por costo (None al final)
orden = sorted(RESULTADOS.items(), key=lambda x: (x[1]['costo'] is None, x[1]['costo'] or 9999))

cols6 = ["Rank","Módulo","Algoritmo","Familia","Mejor Costo\n(días)","vs Línea Base","Tiempo","Evaluaciones","Fuente"]
t6 = s.shapes.add_table(len(orden)+1, len(cols6), Inches(0.2), Inches(1.45), Inches(12.93), Inches(5.75)).table
for i,w in enumerate([0.5,0.6,2.2,1.25,1.1,1.1,0.95,1.1,0.93]):
    t6.columns[i].width = Inches(w)
for c,h in enumerate(cols6):
    celda(t6, 0, c, h, bold=True, size=10, bg=AZUL, fg=BLANCO)

for r,(m,d) in enumerate(orden, 1):
    bg = FAMILIAS_COLOR.get(d["familia"].split("-")[0] if "-" in d["familia"] else d["familia"], RGBColor(0xF9,0xF9,0xF9))
    costo = d["costo"]
    if costo is not None:
        mejora = f"−{round((270.73-costo)/270.73*100,1)}%"
        costo_str = f"{costo:.2f}"
    else:
        mejora = "—"
        costo_str = "—"
    tiempo_str = f"{d['tiempo_h']} h" if d['tiempo_h'] else "—"
    n_eval_str = str(d['n_eval']) if d['n_eval'] else "—"
    rank_str = str(r) if costo is not None else "—"

    fila6 = [rank_str, m, d["algo"], d["familia"].replace("-SK",""), costo_str, mejora, tiempo_str, n_eval_str, d["fuente"]]
    for c,val in enumerate(fila6):
        aln = PP_ALIGN.LEFT if c==2 else PP_ALIGN.CENTER
        fg = VERDE if (costo is not None and costo == 205.17) else (ROJO if costo is None else GRIS_OSCURO)
        if c == 5 and costo is not None:
            fg = VERDE
        celda(t6, r, c, val, size=10, bg=bg, align=aln, fg=fg)

nota6 = s.shapes.add_textbox(Inches(0.2), Inches(7.15), Inches(12.93), Inches(0.28))
set_txt(nota6.text_frame, "★ M4 mejor resultado global  ·  M11/M12 sin datos (colapso de recursos en nube, timeout)  ·  M14 no mejoró línea base (convergencia prematura)", size=8.5, color=GRIS_MEDIO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — INCUMBENTES SMAC (M4–M10)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Parámetros Óptimos (Incumbentes) — Algoritmos Bayesianos M4–M10",
             "Configuración que minimizó tts_full_days_mean en cada módulo")

modulos_bay = ["M4","M7","M8","M9","M10"]
t7 = s.shapes.add_table(len(PARAM_NAMES)+1, len(modulos_bay)+1, Inches(0.2), Inches(1.45), Inches(12.93), Inches(5.75)).table
t7.columns[0].width = Inches(2.43)
for i in range(len(modulos_bay)):
    t7.columns[i+1].width = Inches(2.1)

celda(t7, 0, 0, "Parámetro", bold=True, size=10, bg=AZUL, fg=BLANCO, align=PP_ALIGN.LEFT)
for c,m in enumerate(modulos_bay, 1):
    celda(t7, 0, c, f"{m}\n({RESULTADOS[m]['costo']:.2f} d)", bold=True, size=10, bg=AZUL, fg=BLANCO)

for r,key in enumerate(PARAM_NAMES, 1):
    bg = AZUL_CLARO if r%2==0 else RGBColor(0xFF,0xFF,0xFF)
    celda(t7, r, 0, key, size=9, bg=bg, align=PP_ALIGN.LEFT)
    vals = [INCUMBENTES[m][key] for m in modulos_bay]
    for c,(m,v) in enumerate(zip(modulos_bay,vals), 1):
        val_str = f"{v:.3f}" if isinstance(v,float) else str(v)
        fg = VERDE if ("pct" in key and isinstance(v,float) and v==min(vals)) else GRIS_OSCURO
        celda(t7, r, c, val_str, size=9, bg=bg, fg=fg)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — INCUMBENTES M13/M14 (RAÚL)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Parámetros Óptimos (Incumbentes) — Gradiente Estocástico M13/M14",
             "Resultados de PC Raúl — max_iter=30")

mods_grad = ["M13","M14"]
labels_g  = ["M13 SPSA\n(219.38 días, Raúl PC)", "M14 ALOE\n(270.73 días — sin mejora)"]
t8 = s.shapes.add_table(len(PARAM_NAMES)+1, 3, Inches(1.4), Inches(1.5), Inches(10.1), Inches(5.7)).table
t8.columns[0].width = Inches(3.3); t8.columns[1].width = Inches(3.4); t8.columns[2].width = Inches(3.4)

celda(t8, 0, 0, "Parámetro", bold=True, size=11, bg=AZUL, fg=BLANCO, align=PP_ALIGN.LEFT)
for c,lbl in enumerate(labels_g, 1):
    celda(t8, 0, c, lbl, bold=True, size=11, bg=AZUL, fg=BLANCO)

for r,key in enumerate(PARAM_NAMES, 1):
    bg = AZUL_CLARO if r%2==0 else RGBColor(0xFF,0xFF,0xFF)
    celda(t8, r, 0, key, size=10, bg=bg, align=PP_ALIGN.LEFT)
    for c,m in enumerate(mods_grad, 1):
        v = INCUMBENTES[m][key]
        val_str = f"{v:.3f}" if isinstance(v,float) else str(v)
        fg = NARANJA if (m=="M14" and "pct" in key and isinstance(v,float) and v>0.2) else GRIS_OSCURO
        celda(t8, r, c, val_str, size=10, bg=bg, fg=fg)

nota8 = s.shapes.add_textbox(Inches(0.2), Inches(7.15), Inches(12.93), Inches(0.28))
set_txt(nota8.text_frame, "M14 ALOE: pct_bloqueo=27.5% — convergió a valores altos de bloqueo, no redujo ineficiencia. Algoritmo sensible a ruido elevado y alta dimensión.", size=9, color=NARANJA)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ANÁLISIS INCUMBENTES: PATRONES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Análisis de Incumbentes — Patrones Comunes", "¿Qué variables mueven todos los algoritmos?")

patrones = [
    ("horas_especialista_1ra = 16 h (máximo)",
     "TODOS los módulos con buenos resultados fijan el especialista al máximo.\nNo es el cuello de botella — la eficiencia de esas horas sí lo es."),
    ("pct_bloqueo_1ra → mínimo (~5%)",
     "M4–M10 (Bayesianos) y M13 (SPSA) reducen bloqueo de primera consulta al mínimo.\nEfecto: capacidad efectiva = 16×(1−0.05)×(1−0.05) = 14.4 h vs baseline 7.7 h."),
    ("pct_consultas_vacias → mínimo (~5%)",
     "Misma lógica: eliminar slots vacíos multiplica la capacidad real disponible."),
    ("cupos_lab_ugd = 70–100 (alto)",
     "Laboratorio no es cuello de botella pero requiere holgura para no bloquear flujo."),
    ("M14 ALOE: patrón anómalo",
     "pct_bloqueo = 27.5% en todas las variables — convergencia prematura a un punto\nde simetría artificial. Algoritmo inadecuado para este espacio de alta dimensión y ruido."),
    ("M11/M12 Trust-Region: sin datos",
     "No completaron ninguna iteración tras 4 horas en nube. Requieren ~2h por iteración\n(n réplicas adaptativas ~30 por punto). Pendientes de ejecución en entorno más potente."),
]

y_p = 1.42
for tit_p, cuerpo_p in patrones:
    bar = s.shapes.add_shape(1, Inches(0.35), Inches(y_p+0.1), Inches(0.07), Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = AZUL; bar.line.fill.background()
    bx = s.shapes.add_textbox(Inches(0.55), Inches(y_p), Inches(12.43), Inches(0.9))
    tf = bx.text_frame; tf.word_wrap = True
    set_txt(tf, tit_p, bold=True, size=11, color=AZUL)
    add_txt(tf, cuerpo_p, size=10, color=GRIS_OSCURO)
    y_p += 0.97

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — GRÁFICAS INDIVIDUALES NUBE (4 módulos)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Gráficas de Convergencia — Módulos Bayesianos Nube (M4, M8, M9, M10)")

PNGS4 = [
    ("M4",  "/home/user/SO/resultado_comparativa_m4.png"),
    ("M8",  "/home/user/SO/resultado_comparativa_m8.png"),
    ("M9",  "/home/user/SO/resultado_comparativa_m9.png"),
    ("M10", "/home/user/SO/resultado_comparativa_m10.png"),
]
for (m,path),(x,y,w,h) in zip(PNGS4, [(0.25,1.45,6.3,2.95),(6.75,1.45,6.3,2.95),(0.25,4.45,6.3,2.95),(6.75,4.45,6.3,2.95)]):
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        bx = s.shapes.add_textbox(Inches(x), Inches(y+1.2), Inches(w), Inches(0.5))
        set_txt(bx.text_frame, f"[{m}: imagen no disponible]", size=10, color=NARANJA, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — GRÁFICAS RAÚL (M13, M14)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Gráficas de Convergencia — M13 SPSA y M14 ALOE (PC Raúl, max_iter=30)")

PNGS_RAUL = [
    ("M13 SPSA (219.38 d)", "/tmp/raul_v2/Raul_23_05_2026/resultado_comparativa_m13.png"),
    ("M14 ALOE (270.73 d)", "/tmp/raul_v2/Raul_23_05_2026/resultado_comparativa_m14.png"),
]
for (lbl,path),(x,y,w,h) in zip(PNGS_RAUL, [(0.4,1.45,6.1,5.7),(6.83,1.45,6.1,5.7)]):
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — COMPARATIVA CONVERGENCIA (RAÚL)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Comparativa de Convergencia — PC Raúl (M4 vs M13 vs M14)")
p = "/tmp/raul_v2/Raul_23_05_2026/comparativa_convergencia.png"
if os.path.exists(p):
    s.shapes.add_picture(p, Inches(0.5), Inches(1.48), Inches(12.33), Inches(5.75))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — COMPARATIVA CONFIGURACIONES (RAÚL)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Comparativa de Configuraciones Óptimas — PC Raúl (M4 vs M13 vs M14)")
p = "/tmp/raul_v2/Raul_23_05_2026/comparativa_configuraciones.png"
if os.path.exists(p):
    s.shapes.add_picture(p, Inches(0.5), Inches(1.48), Inches(12.33), Inches(5.75))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — TABLA RESUMEN (RAÚL)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Tabla Resumen — PC Raúl (M4, M13, M14)")
p = "/tmp/raul_v2/Raul_23_05_2026/comparativa_tabla_resumen.png"
if os.path.exists(p):
    s.shapes.add_picture(p, Inches(0.5), Inches(1.48), Inches(12.33), Inches(5.75))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — M10 HPO KGCP (RESULTADO ADICIONAL RAÚL)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Resultado Adicional — M10 HPO KGCP (PC Raúl)",
             "Modo HPO con 100 evaluaciones: segundo mejor resultado global")

hpo_data = [
    ("Algoritmo",          "SK-KGCP en modo HPO (Knowledge Gradient + Correlated Prior)"),
    ("Costo incumbente",   "212.03 días  ★  (segundo mejor resultado, solo superado por M4=205.17)"),
    ("Tiempo total",       "12.27 horas"),
    ("Evaluaciones",       "100 trials (vs 20 del modo estándar)"),
    ("Diferencia vs M10 estándar (nube)", "M10 std: 219.96 días (20 eval, 1.46h)  →  HPO: 212.03 días (100 eval, 12.27h)\nMás evaluaciones permiten al surrogado SK-KGCP mostrar su ventaja real."),
    ("Incumbente", "cupos_eco_mat=19, cupos_eco_ugd=21, cupos_lab=81, dias_pub=8\nhoras_control=50, h_esp_1ra=22, agentes=2, matronas=2\npct_bloqueo_1ra=0.151, pct_bloqueo_post=0.212, pct_vacias=0.215, pct_no_cont=0.446"),
]

y_hpo = 1.44
for label, valor in hpo_data:
    bar = s.shapes.add_shape(1, Inches(0.35), Inches(y_hpo+0.08), Inches(0.07), Inches(0.6))
    bar.fill.solid(); bar.fill.fore_color.rgb = AZUL; bar.line.fill.background()
    bx = s.shapes.add_textbox(Inches(0.55), Inches(y_hpo), Inches(12.43), Inches(0.92))
    tf = bx.text_frame; tf.word_wrap = True
    set_txt(tf, label, bold=True, size=11, color=AZUL)
    fg_val = VERDE if "★" in valor else GRIS_OSCURO
    add_txt(tf, valor, size=10.5, color=fg_val)
    y_hpo += 0.97

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — CONVERGENCIA DETALLADA M13 SPSA
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Convergencia Detallada — M13 SPSA (PC Raúl, max_iter=30)",
             "Historial completo: 31 evaluaciones incumbentes")

# tabla de iteraciones clave
iters_key = [
    (0,  270.73, 270.73, "—",     "Punto inicial (centro espacio)"),
    (1,  235.39, 235.39, "0.39h", "Primera mejora: −13% de golpe"),
    (10, 222.03, 222.03, "2.75h", "Segunda mejora significativa"),
    (13, 224.67, 222.03, "3.60h", "Evaluación ruidosa, incumbente mantiene"),
    (21, 219.38, 219.38, "6.28h", "MEJOR ENCONTRADO: 219.38 días"),
    (30, 242.74, 219.38, "8.65h", "Fin: evaluación ruidosa, incumbente 219.38"),
]

cols_s = ["Iter","Costo eval","Incumbente","Tiempo","Observación"]
t_s = s.shapes.add_table(len(iters_key)+1, 5, Inches(0.3), Inches(1.48), Inches(12.73), Inches(3.5)).table
for i,w in enumerate([0.6,1.2,1.3,0.9,8.0]):
    t_s.columns[i].width = Inches(w)
for c,h in enumerate(cols_s):
    celda(t_s, 0, c, h, bold=True, size=11, bg=AZUL, fg=BLANCO)

BG_ALT = [RGBColor(0xFF,0xFF,0xFF), AZUL_CLARO]
for r,(it,ceval,cinc,t,obs) in enumerate(iters_key, 1):
    bg = BG_ALT[r%2]
    fg_eval = VERDE if ceval == cinc and ceval < 270 else (NARANJA if ceval > 290 else GRIS_OSCURO)
    fg_inc  = VERDE if cinc == 219.38 else GRIS_OSCURO
    celda(t_s, r, 0, it,    size=11, bg=bg)
    celda(t_s, r, 1, f"{ceval:.2f}", size=11, bg=bg, fg=fg_eval)
    celda(t_s, r, 2, f"{cinc:.2f}",  size=11, bg=bg, fg=fg_inc)
    celda(t_s, r, 3, t,     size=11, bg=bg)
    celda(t_s, r, 4, obs,   size=11, bg=bg, align=PP_ALIGN.LEFT)

bx_obs = s.shapes.add_textbox(Inches(0.3), Inches(5.1), Inches(12.73), Inches(2.1))
tf_obs = bx_obs.text_frame; tf_obs.word_wrap = True
set_txt(tf_obs, "Análisis de convergencia SPSA:", bold=True, size=12, color=AZUL)
add_txt(tf_obs, "• Alta varianza entre evaluaciones (rango: 219–320 días) — ruido del simulador domina la señal del gradiente", size=11, color=GRIS_OSCURO)
add_txt(tf_obs, "• El incumbente se actualiza solo en iteraciones 1, 10, 21 (3 de 30) — exploración ineficiente", size=11, color=GRIS_OSCURO)
add_txt(tf_obs, "• Tendencia decreciente del paso αk (0.142 → 0.065): el algoritmo reduce exploración con el tiempo", size=11, color=GRIS_OSCURO)
add_txt(tf_obs, "• Resultado final: 219.38 días — mejora del 19% sobre línea base, en 8.65 h", size=11, color=VERDE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — DIAGNÓSTICO M14 ALOE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Diagnóstico — M14 ALOE: Por qué no mejoró (PC Raúl, max_iter=30)",
             "armijo_ok=False en las 30 iteraciones — Armijo rechazó TODOS los pasos")

diag = [
    ("Síntoma observado",
     "tts_full_days_mean = 270.73 días en TODAS las 30 iteraciones.\nEl incumbente nunca mejoró respecto al punto inicial."),
    ("Causa: Armijo siempre rechazado",
     "ALOE usa búsqueda de línea de Armijo para aceptar/rechazar cada paso.\nCondición: f(x + α·d) ≤ f(x) − c·α·‖d‖²\nEn las 30 iteraciones, armijo_ok=False SIEMPRE → α se redujo de 1.0 a 0.0012."),
    ("Por qué falla Armijo aquí",
     "El simulador DES es altamente estocástico (σ ≈ 20–40 días).\nCuando el ruido supera la mejora esperada del paso, Armijo lo rechaza por precaución.\nCon 12 variables y ruido alto, la dirección de gradiente estimada es poco confiable."),
    ("Comparación: SPSA vs ALOE",
     "SPSA acepta pasos aunque sean ruidosos → explora más, encuentra 219 días.\nALOE es más conservador (Armijo) → queda atrapado en el punto inicial.\nConclusión: en simuladores ruidosos de alta dimensión, SPSA supera a ALOE."),
    ("Recomendación",
     "ALOE es adecuado para funciones deterministas o bajo ruido (σ < 5%).\nPara este simulador (σ ≈ 15%), usar SPSA, ASTRO-DF, o métodos Bayesianos."),
]

y_diag = 1.44
for tit_d, cuerpo_d in diag:
    bar = s.shapes.add_shape(1, Inches(0.35), Inches(y_diag+0.08), Inches(0.07), Inches(0.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = ROJO; bar.line.fill.background()
    bx = s.shapes.add_textbox(Inches(0.55), Inches(y_diag), Inches(12.43), Inches(0.9))
    tf = bx.text_frame; tf.word_wrap = True
    set_txt(tf, tit_d, bold=True, size=11, color=ROJO)
    add_txt(tf, cuerpo_d, size=10, color=GRIS_OSCURO)
    y_diag += 1.18 if len(cuerpo_d) > 80 else 0.98

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — MÓDULOS PENDIENTES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s, AMARILLO_BG)
banda_titulo(s, "Módulos Pendientes de Ejecución — M11 y M12",
             "Trust-Region: ASTRO-DF y STRONG requieren entorno con mayor presupuesto de tiempo")

pendientes = [
    ("M11 — ASTRO-DF", "Poloczek et al. 2017",
     "Evaluación de x0: ~28 min (n≈30 réplicas adaptativas)\nCada iteración: ~28–40 min → 3 iteraciones ≈ 1.5 h mínimo\nMotivo de fallo: pool de workers agotó RAM tras 4h de ejecución continua"),
    ("M12 — STRONG",   "Chang et al. 2013",
     "Evaluación de x0: >2 h (n0=10 puntos × n_r=10 réplicas = 100 evaluaciones)\nCada iteración: ~60–90 min → 3 iteraciones ≈ 4–5 h total\nMotivo de fallo: nunca completó x0 en la sesión de nube"),
]

y_pend = 1.5
for titulo_pend, ref_pend, detalle_pend in pendientes:
    rect = s.shapes.add_shape(1, Inches(0.4), Inches(y_pend), Inches(12.5), Inches(2.3))
    rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF)
    rect.line.color.rgb = AZUL
    bx = s.shapes.add_textbox(Inches(0.6), Inches(y_pend+0.1), Inches(12.1), Inches(2.0))
    tf = bx.text_frame; tf.word_wrap = True
    set_txt(tf, f"{titulo_pend}  ({ref_pend})", bold=True, size=13, color=AZUL)
    add_txt(tf, detalle_pend, size=11, color=GRIS_OSCURO)
    y_pend += 2.5

bx_rec = s.shapes.add_textbox(Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.75))
tf_rec = bx_rec.text_frame; tf_rec.word_wrap = True
set_txt(tf_rec, "Recomendación: ejecutar en GitHub Actions (hasta 6h) o PC local con continuar_comparativa.py", bold=True, size=12, color=NARANJA)
add_txt(tf_rec, "El script detecta M4–M10, M13, M14 ya completados y corre solo M11 y M12.", size=11, color=GRIS_OSCURO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — CONCLUSIONES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
fondo(s)
banda_titulo(s, "Conclusiones")

conclusiones = [
    ("Mejor resultado: M4 SMAC GP+EI (205.17 días, −24%)",
     "Con solo 20 evaluaciones y ~1.3 h, el algoritmo más simple entregó la mejor solución.\nSurrogado GP+EI es el más adecuado para 12 variables con ruido alto."),
    ("Algoritmos Bayesianos (M4–M10): todos convergen",
     "Rango 205–220 días. Las variantes SK (M7–M10) no superan a GP básico (M4),\nsugeriendo que el ruido del simulador limita la ventaja de modelos más sofisticados."),
    ("Gradiente estocástico: SPSA funciona, ALOE no",
     "M13 SPSA: 219.38 días con 30 iteraciones (~8.65 h) — solución razonable.\nM14 ALOE: 270.73 días — sin mejora tras 69.5 h. Inadecuado para alta dimensión y ruido."),
    ("Trust-Region (M11/M12): pendientes",
     "Requieren 4–6 h por ejecución. No completaron en entorno nube (timeout + OOM).\nPendientes de ejecutar en PC local o GitHub Actions."),
    ("Hallazgo clave: eficiencia de slots > dotación",
     "La palanca real es pct_bloqueo_1ra y pct_consultas_vacias (0.05 vs 0.30 baseline).\nReducirlos de 30% a 5% equivale a duplicar capacidad efectiva del especialista."),
]

y_c = 1.44
for tit_c, cuerpo_c in conclusiones:
    bar = s.shapes.add_shape(1, Inches(0.35), Inches(y_c+0.08), Inches(0.07), Inches(0.6))
    bar.fill.solid(); bar.fill.fore_color.rgb = AZUL; bar.line.fill.background()
    bx = s.shapes.add_textbox(Inches(0.55), Inches(y_c), Inches(12.43), Inches(1.0))
    tf = bx.text_frame; tf.word_wrap = True
    set_txt(tf, tit_c, bold=True, size=11, color=AZUL)
    add_txt(tf, cuerpo_c, size=10, color=GRIS_OSCURO)
    y_c += 1.02

# ── Guardar ───────────────────────────────────────────────────────────────────
out = "/home/user/SO/comparativa_optimizadores_v2.pptx"
prs.save(out)
print(f"PPT v2 guardado: {out}")
